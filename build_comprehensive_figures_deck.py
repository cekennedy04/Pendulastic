"""
build_comprehensive_figures_deck.py
====================================
"Every figure we can think of" deck (2026-08-20), in the same visual system as
the Pendulastic Progress Update.pptx palette (Georgia/Times New Roman, navy-on-
off-white, stat cards, key-finding cards).

Revision (this pass): fixed real layout bugs from the first cut -- several
figures were being placed at a fixed pixel width with no height cap, so a
tall/narrow chart (e.g. the data-availability matrix, aspect ratio 0.81) or a
near-square heatmap (aspect ratio 0.94) overflowed past the bottom of the
slide. Every figure is now placed via place_in_box(), which reads the PNG's
real aspect ratio and fits it inside a fixed content box (contain, not
stretch), centered. Dense multi-panel figures (the 7-parameter small
multiples, the scorecard, the group-comparison charts) got their own full-
width slide instead of being squeezed to half width and becoming illegible.
Headers and dividers were also simplified and made consistent -- one header
treatment for content slides, Georgia reserved for section breaks and the
title/closing only, and the repetitive "Figures N, M | ..." caption formula
was dropped in favor of plain, varied editorial lines (or none, when the
figure speaks for itself).

Usage:
    .venv\\Scripts\\python.exe build_comprehensive_figures_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

FIG_DIR = "Model_Analysis_Outputs/paper_figures"
OUT_PATH = "Model_Analysis_Outputs/Pendulastic_All_Figures_Deck_v4.pptx"

BG = RGBColor(0xF0, 0xF5, 0xFA)
NAVY = RGBColor(0x0C, 0x1E, 0x34)
SLATE = RGBColor(0x3A, 0x58, 0x70)
GRAY_HEAD = RGBColor(0x60, 0x62, 0x64)
BODY = RGBColor(0x0A, 0x0A, 0x0A)
STAT_NUM = RGBColor(0x16, 0x20, 0x2C)
STAT_DESC = RGBColor(0x4A, 0x55, 0x68)
CARD_BORDER = RGBColor(0xB8, 0xCD, 0xDE)
NUM_LIGHT = RGBColor(0xB8, 0xCD, 0xDE)
RULE = RGBColor(0xD8, 0xDE, 0xE6)
ACCENT_PURPLE = RGBColor(0x60, 0x30, 0xA0)
ACCENT_AMBER = RGBColor(0xB0, 0x60, 0x10)
ACCENT_GREEN = RGBColor(0x0A, 0x7B, 0x52)
ACCENT_RED = RGBColor(0xC0, 0x30, 0x30)
ACCENT_TEAL = RGBColor(0x1E, 0x4F, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Georgia"
BODY_FONT = "Times New Roman"

# ---- layout grid (inches) ----
PAGE_W, PAGE_H = 13.333, 7.5
MARGIN_L, MARGIN_R = 0.85, 0.85
CONTENT_W = PAGE_W - MARGIN_L - MARGIN_R
TOP_NO_CAPTION = 1.55
TOP_WITH_CAPTION = 1.95
BOTTOM = 7.05
GUTTER = 0.5

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def tb(slide, l, t, w, h, text, size, color, font=BODY_FONT, bold=False, italic=False,
       align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = box.text_frame.paragraphs[0] if i == 0 else box.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = font
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = align
    return box


def header(slide, title, caption=None):
    """One consistent header treatment for every content slide."""
    tb(slide, MARGIN_L, 0.55, CONTENT_W, 0.55, title, 27, GRAY_HEAD, font=BODY_FONT, bold=False)
    ln = slide.shapes.add_connector(1, Inches(MARGIN_L), Inches(1.14), Inches(MARGIN_L + CONTENT_W), Inches(1.14))
    ln.line.color.rgb = RULE
    ln.line.width = Pt(1)
    if caption:
        tb(slide, MARGIN_L, 1.28, CONTENT_W, 0.4, caption, 13, STAT_DESC, font=BODY_FONT)
        return TOP_WITH_CAPTION
    return TOP_NO_CAPTION


def divider_slide(title, kicker):
    s = new_slide()
    tb(s, 0.0, 3.05, PAGE_W, 0.4, kicker.upper(), 12.5, SLATE, font=BODY_FONT, align=PP_ALIGN.CENTER)
    tb(s, 1.0, 3.5, PAGE_W - 2.0, 1.2, title, 34, NAVY, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    ln = slide_rule = s.shapes.add_connector(1, Inches(PAGE_W / 2 - 0.6), Inches(4.35), Inches(PAGE_W / 2 + 0.6), Inches(4.35))
    ln.line.color.rgb = RULE
    ln.line.width = Pt(1.25)
    return s


_img_size_cache = {}


def img_size(path):
    if path not in _img_size_cache:
        with Image.open(path) as im:
            _img_size_cache[path] = im.size
    return _img_size_cache[path]


def place_in_box(slide, filename, box_l, box_t, box_w, box_h, valign="middle"):
    """Fit an image inside a box, preserving aspect ratio, centered (contain, never stretch)."""
    path = os.path.join(FIG_DIR, filename)
    if not os.path.exists(path):
        tb(slide, box_l, box_t, box_w, 0.4, f"[missing: {filename}]", 11, BODY)
        return
    iw, ih = img_size(path)
    ratio = iw / ih
    w, h = box_w, box_w / ratio
    if h > box_h:
        h = box_h
        w = box_h * ratio
    left = box_l + (box_w - w) / 2
    if valign == "top":
        top = box_t
    else:
        top = box_t + (box_h - h) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(w), height=Inches(h))


def full_figure(slide, filename, top=None, bottom=BOTTOM, caption=None):
    top = top if top is not None else TOP_NO_CAPTION
    box_h = bottom - top - (0.4 if caption else 0)
    place_in_box(slide, filename, MARGIN_L, top, CONTENT_W, box_h)
    if caption:
        tb(slide, MARGIN_L, bottom - 0.35, CONTENT_W, 0.4, caption, 11.5, STAT_DESC, font=BODY_FONT, italic=True)


def two_up(slide, fname_l, fname_r, top=None, bottom=BOTTOM, caption=None):
    top = top if top is not None else TOP_NO_CAPTION
    box_h = bottom - top - (0.4 if caption else 0)
    box_w = (CONTENT_W - GUTTER) / 2
    place_in_box(slide, fname_l, MARGIN_L, top, box_w, box_h)
    place_in_box(slide, fname_r, MARGIN_L + box_w + GUTTER, top, box_w, box_h)
    if caption:
        tb(slide, MARGIN_L, bottom - 0.35, CONTENT_W, 0.4, caption, 11.5, STAT_DESC, font=BODY_FONT, italic=True)


def add_table(slide, headers_, rows, left, top, width, height, font_size=11, header_fill=NAVY):
    n_rows = len(rows) + 1
    n_cols = len(headers_)
    tshape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = tshape.table
    for c, h in enumerate(headers_):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.name = BODY_FONT
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].font.name = BODY_FONT
            cell.text_frame.paragraphs[0].font.color.rgb = BODY
    return table


def stat_card(slide, left, top, width, number, desc):
    tb(slide, left, top, width, 0.5, number, 23, STAT_NUM, font=BODY_FONT, bold=True)
    tb(slide, left, top + 0.44, width, 0.6, desc, 10.5, STAT_DESC, font=BODY_FONT)


def key_finding_card(slide, top, num, heading, body, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(top), Inches(CONTENT_W), Inches(1.0))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CARD_BORDER; card.line.width = Pt(0.75)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(top), Inches(0.07), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tb(slide, MARGIN_L + 0.2, top + 0.15, 0.5, 0.6, str(num), 23, NUM_LIGHT, font=TITLE_FONT)
    tb(slide, MARGIN_L + 1.75, top + 0.08, CONTENT_W - 2.0, 0.38, heading, 14.5, NAVY, font=TITLE_FONT)
    tb(slide, MARGIN_L + 1.75, top + 0.46, CONTENT_W - 2.0, 0.46, body, 10.5, SLATE, font="Calibri")


def roadmap_row(slide, top, phase, desc, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MARGIN_L + 0.05), Inches(top + 0.06), Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    tb(slide, MARGIN_L + 0.5, top, 1.7, 0.4, phase, 14, NAVY, font=TITLE_FONT, bold=True)
    tb(slide, MARGIN_L + 2.3, top - 0.02, CONTENT_W - 2.3, 0.65, desc, 12, BODY, font=BODY_FONT)


# ================================================================= Title
s = new_slide()
tb(s, 0.0, 2.5, PAGE_W, 1.1, "Pendulastic", 46, NAVY, font=TITLE_FONT, align=PP_ALIGN.CENTER)
tb(s, 0.0, 3.55, PAGE_W, 0.55, "Comprehensive Results & Figures", 22, SLATE, font=TITLE_FONT, align=PP_ALIGN.CENTER)
tb(s, 0.0, 6.85, PAGE_W, 0.4, "30 figures across instrument validation, calibration methodology, "
   "clinical correlation, and next steps  --  2026-08-20", 12.5, STAT_DESC, font=BODY_FONT, align=PP_ALIGN.CENTER)

# ================================================================= Part 1: Instrument Validation
divider_slide("Instrument Validation", "Part 1")

s = new_slide()
top = header(s, "Where does the IMU's error come from?", "53 trials, full-curve RMSE against OptiTrack")
two_up(s, "fig11_rmse_distribution.png", "fig13_bias_vs_rmse.png", top=top,
      caption="Left: the spread of RMSE across trials. Right: bias alone explains almost all of it (R² shown on chart).")

s = new_slide()
top = header(s, "Two confounds, ruled out")
two_up(s, "fig14_lag_distribution.png", "fig23_rmse_vs_trial_length.png", top=top,
      caption="Sync-lag drift and clip length were both checked as alternative explanations for the error -- neither holds.")

s = new_slide()
top = header(s, "How the error breaks down by participant")
two_up(s, "fig12_rmse_by_participant.png", "fig1_bland_altman.png", top=top,
      caption="No single participant drives the overall error level (left); agreement on the relaxation index specifically (right).")

s = new_slide()
top = header(s, "Table 1. Agreement per PT parameter", "IMU vs. OptiTrack, n = 61 trials, 5 participants, ICC(2,1)")
add_table(s, ["Parameter", "ICC(2,1)", "Bias", "95% LoA"], [
    ["R2n (relaxation index)", "0.226", "-0.108", "[-1.06, 0.85]"],
    ["N (oscillation count)", "0.458", "-1.46", "[-7.48, 4.56]"],
    ["phi_max ratio", "0.044", "-0.055", "[-0.51, 0.40]"],
    ["omega_max,n", "0.014", "-2.13", "[-12.05, 7.78]"],
    ["f (frequency)", "0.140", "-0.40", "[-1.98, 1.18]"],
    ["Area ratio", "0.135", "-0.052", "[-0.82, 0.72]"],
    ["omega_min,n", "0.214", "-1.01", "[-6.70, 4.69]"],
], left=1.7, top=1.95, width=9.9, height=3.35, font_size=13.5)
stat_card(s, 1.7, 5.65, 3.2, "14.84 / 10.98°", "Full-curve RMSE, mean / median")
stat_card(s, 5.55, 5.65, 3.1, "2 of 53", "Trials met the 5° clinical-goal threshold")
stat_card(s, 9.1, 5.65, 2.6, "67-76%", "Of total error is fixed calibration bias")

# ================================================================= Part 2: Calibration Methodology
divider_slide("The Calibration Methodology Search", "Part 2")

s = new_slide()
top = header(s, "A 144-combination grid search", "beta × ema_alpha × flex_axis_capture × gravity_seed × method")
full_figure(s, "fig22_grid_search_convergence.png", top=top,
           caption="The adopted configuration sits at the global minimum of the full sweep, not a local one.")

s = new_slide()
top = header(s, "Which model, and how much did tuning matter?")
two_up(s, "fig17_method_comparison.png", "fig18_beta_sensitivity.png", top=top,
      caption="The win came from picking the right model, not from fine-tuning beta -- RMSE is nearly flat across it.")

s = new_slide()
top = header(s, "Two dead ends, investigated and closed out")
two_up(s, "fig15_ft_ratio_sweep.png", "fig16_mag_toggle_paired.png", top=top,
      caption="ft_ratio never converges to a real optimum; magnetometer fusion changes nothing.")

# ================================================================= Part 3: Markerless Video
divider_slide("Markerless Video (MediaPipe)", "Part 3")

s = new_slide()
top = header(s, "MediaPipe against the same OptiTrack reference")
box_w = (CONTENT_W - GUTTER) / 2
place_in_box(s, "fig31_mp_trajectory_example.png", MARGIN_L, top, box_w, BOTTOM - top - 0.55)
add_table(s, ["Parameter", "ICC(2,1)", "Bias", "95% LoA"], [
    ["R2n", "-0.041", "-0.441", "[-2.35, 1.47]"],
    ["N", "0.032", "-2.30", "[-18.35, 13.75]"],
    ["phi_max ratio", "-0.008", "0.110", "[-0.91, 1.13]"],
    ["omega_max,n", "-0.036", "6.75", "[-50.28, 63.77]"],
    ["f", "-0.115", "-0.83", "[-3.04, 1.37]"],
    ["Area ratio", "0.003", "-0.420", "[-1.23, 0.39]"],
    ["omega_min,n", "-0.018", "3.69", "[-29.16, 36.53]"],
], left=MARGIN_L + box_w + GUTTER, top=top, width=box_w, height=3.15, font_size=10.5)
tb(s, MARGIN_L + box_w + GUTTER, top + 3.35, box_w, 1.3,
   "Full-curve RMSE 36.0° mean / 33.3° median (n=37). ICC is at or below zero for every parameter -- "
   "no measurable agreement, well behind the IMU.", 12, BODY, font=BODY_FONT)

# ================================================================= Part 4: Cohort & Clinical
divider_slide("Cohort, Clinical Grading & Data Availability", "Part 4")

s = new_slide()
top = header(s, "Who's in the dataset, and what do they have?")
two_up(s, "fig19_mas_distribution.png", "fig20_data_availability.png", top=top,
      caption="The right-hand matrix is the root cause behind most of the limitations later in this deck.")

s = new_slide()
top = header(s, "PT metrics by clinical category", "format matches Whelan et al. (2018), Figure 3")
full_figure(s, "fig2_metrics_by_mas.png", top=top)

# ================================================================= Part 5: Group Differences & Multi-modal
divider_slide("Group Differences & Multi-Modal Agreement", "Part 5")

s = new_slide()
top = header(s, "Control vs. MS, headline metrics")
full_figure(s, "fig4_metrics_by_group.png", top=top,
           caption="Each open circle is one participant's own mean -- the control arm is visibly one person, not a distribution.")

s = new_slide()
top = header(s, "All 7 parameters, IMU: Control vs. MS")
full_figure(s, "fig21_pt_params_small_multiples.png", top=top)

s = new_slide()
top = header(s, "The same comparison, OptiTrack ground truth")
full_figure(s, "fig24_pt_params_small_multiples_opti.png", top=top,
           caption="Cleanest Control/MS separation of the three modalities -- worth holding up against the IMU and MediaPipe panels.")

s = new_slide()
top = header(s, "The same comparison, MediaPipe")
full_figure(s, "fig25_pt_params_small_multiples_mp.png", top=top,
           caption="Boxes largely overlap -- consistent with MediaPipe's near-zero ICC on Table 2.")

s = new_slide()
top = header(s, "Do the metrics even agree with each other?")
two_up(s, "fig6_metric_effect_heatmap.png", "fig7_single_vs_combined_auc.png", top=top,
      caption="Left: sign of the effect flips by modality for most parameters. Right: the honest generalization test -- every AUC below chance.")

# ================================================================= Part 6: Composite Score Critique
divider_slide("Composite PT Score, a Critical Look", "Part 6")

s = new_slide()
top = header(s, "Does the production score separate groups, and is it sound?")
two_up(s, "fig8_score_naive_vs_logocv.png", "fig10_param_correlation.png", top=top,
      caption="Left: p=0.5865 naively, worse under cross-validation. Right: two of the seven parameters are near-redundant (r=0.93).")

# ================================================================= Part 7: Longitudinal Change
divider_slide("Longitudinal Change", "Part 7")

s = new_slide()
top = header(s, "Pre/post, the one paired observation we have")
full_figure(s, "fig5_pre_post.png", top=top, caption="Illustrative only -- n=1 paired participant, not a statistical comparison.")

s = new_slide()
top = header(s, "Every participant's change over time, all 7 parameters")
full_figure(s, "fig26_longitudinal_all_params.png", top=top,
           caption="Only P15 (both legs) has 2+ matched timepoints in the OptiTrack-validated set. MAS labels are each leg's "
                   "most recently recorded grade, not necessarily the grade at every point plotted.")

s = new_slide()
top = header(s, "R2n over time, in detail")
box_w = 8.6
place_in_box(s, "fig27_r2n_longitudinal_detail.png", MARGIN_L + (CONTENT_W - box_w) / 2, top, box_w, BOTTOM - top)

# ================================================================= Part 8: Composite Score Across Modalities
divider_slide("One Score, Three Modalities", "Part 8")

s = new_slide()
top = header(s, "A session-computed score: OptiTrack vs. IMU vs. MediaPipe",
            "Not the production compute_pt_score() -- that's IMU-only (Figure 8). Each modality here is scored "
            "against its own control median, so all three are compared on equal footing.")
two_up(s, "fig28_score_by_modality.png", "fig29_score_correlation.png", top=top)

s = new_slide()
top = header(s, "Every participant, every leg, every modality, next to their MAS grade")
full_figure(s, "fig30_mas_scorecard.png", top=top)

# ================================================================= Next Steps
s = new_slide()
top = header(s, "Next steps")
roadmap_row(s, top + 0.15, "Step 1", "IMU-record the 7 existing video-only healthy controls -- zero new recruitment, "
                              "resolves the n=1 control-arm bottleneck behind Figures 4, 7, and 10 at once.", ACCENT_GREEN)
roadmap_row(s, top + 0.9, "Step 2", "Field-test the calibration-bias fix on a controlled recording session -- targets the "
                              "67-76% of RMSE that is bias, not noise (Figure 13).", ACCENT_TEAL)
roadmap_row(s, top + 1.65, "Step 3", "Recruit across the full MAS severity range -- current cohort has zero trials at "
                              "MAS ≥ 2 (Figure 19); no severity claim is possible until this exists.", ACCENT_PURPLE)
roadmap_row(s, top + 2.4, "Step 4", "Record mas_extension, not just the collapsed mas_grade, across the full cohort, and "
                              "switch the clinical-correlation target to it -- the pendulum test is extensor-specific.", ACCENT_AMBER)
roadmap_row(s, top + 3.15, "Step 5", "Fix the score formula's internal issues independent of more data: decorrelate or "
                              "reweight the 7 parameters (R2n & omega_max,n at r=0.93) instead of equal-weighting.", ACCENT_AMBER)
roadmap_row(s, top + 3.9, "Step 6", "Re-run the multi-metric AUC test and the naive/LOGO-CV score comparison once Steps "
                              "1-3 land -- only then is a combined-diagnostic claim defensible.", ACCENT_RED)
tb(s, MARGIN_L, top + 4.75, CONTENT_W, 0.5,
   "Ordered by leverage -- Step 1 is free and unblocks three separate figures at once.",
   12, STAT_DESC, font=BODY_FONT, italic=True)

# ================================================================= Key Findings
s = new_slide()
top = header(s, "Key findings, in one place")
key_finding_card(s, top + 0.1, 1, "IMU beats markerless video, with a fixable error source",
                  "ICC 0.01-0.46 vs. MediaPipe's ICC at or below zero on every parameter; roughly 70% of IMU "
                  "error is a correctable per-trial calibration bias.", ACCENT_GREEN)
key_finding_card(s, top + 1.2, 2, "The methodology search converged on a simpler model, not the fancier one",
                  "The 144-combo grid search: the simple relative-quaternion method beats both Ockendon "
                  "variants by 2-5x; ft_ratio and magnetometer fusion were both dead ends.", ACCENT_TEAL)
key_finding_card(s, top + 2.3, 3, "The clinical correlation is real, but the cohort is thin",
                  "R2n vs. MAS: rho=-0.313, p=0.014 -- significant, but about half the strength of the nearest "
                  "published comparator, and zero trials exist at MAS ≥ 2.", ACCENT_PURPLE)
key_finding_card(s, top + 3.4, 4, "Neither single metrics nor the composite score separate groups yet",
                  "Every leave-one-out AUC sits below chance; a naive comparison on the production score gives "
                  "p=0.5865, contradicting prior internal documentation of p=0.0001.", ACCENT_RED)
key_finding_card(s, top + 4.5, 5, "Almost every limitation traces back to one addressable cause",
                  "Only 5 of 14 enrolled participants have IMU data, and the control arm is a single person -- "
                  "see Next Steps, Step 1.", ACCENT_AMBER)

# ================================================================= Closing
s = new_slide()
tb(s, 0.0, 3.1, PAGE_W, 0.8, "Thank you", 40, NAVY, font=TITLE_FONT, align=PP_ALIGN.CENTER)
tb(s, 0.0, 4.0, PAGE_W, 0.5, "30 figures, one dataset, six next steps", 16, SLATE,
   font=TITLE_FONT, align=PP_ALIGN.CENTER)
tb(s, 0.0, 6.85, PAGE_W, 0.4, "Pendulastic  --  2026-08-20", 12, STAT_DESC, font=BODY_FONT, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH}  ({len(prs.slides._sldIdLst)} slides)")
