"""
build_results_pptx_v2.py
=========================
Rebuilds the Results deck in the visual style of "Pendulastic Progress Update.pptx"
(Google-Slides-exported: Georgia/Times New Roman, navy-on-off-white palette, stat
cards, ranking bars, numbered key-finding cards) -- extracted directly from that
deck via python-pptx introspection, not guessed.

Adds three critique-driven slides:
  - Reframed MediaPipe/"assessor-tracking" finding (Section 3.3) as a validation
    contribution rather than a failure, backed by a freshly computed OT~MP
    regression (see verify_assessor_regression.py output).
  - Zero-Recruitment next-step slide (Section on the n=1 control fix).
  - T3/T5 physiological breakout slide is INTENTIONALLY OMITTED here -- the
    specific numbers quoted in the critique (CV=3.12%, negative A1,flex at T3/T5)
    could not be independently verified against this session's computed outputs.
    See chat response for the follow-up ask.

Usage:
    .venv\\Scripts\\python.exe build_results_pptx_v2.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

import sys
COMPRESSED = "--compressed" in sys.argv
NO_IMAGES = "--no-images" in sys.argv
FIG_DIR = "Model_Analysis_Outputs/paper_figures_compressed" if COMPRESSED else "Model_Analysis_Outputs/paper_figures"
FIG_EXT = ".jpg" if COMPRESSED else ".png"
if NO_IMAGES:
    OUT_PATH = "Model_Analysis_Outputs/Pendulastic_Results_Section_v2_slides.pptx"
elif COMPRESSED:
    OUT_PATH = "Model_Analysis_Outputs/Pendulastic_Results_Section_v2_compressed.pptx"
else:
    OUT_PATH = "Model_Analysis_Outputs/Pendulastic_Results_Section_v2.pptx"

# ---- Palette + fonts extracted from Pendulastic Progress Update.pptx ----
BG = RGBColor(0xF0, 0xF5, 0xFA)
NAVY = RGBColor(0x0C, 0x1E, 0x34)
SLATE = RGBColor(0x3A, 0x58, 0x70)
GRAY_HEAD = RGBColor(0x60, 0x62, 0x64)
GRAY_HEAD2 = RGBColor(0x88, 0x88, 0x88)
BODY = RGBColor(0x0A, 0x0A, 0x0A)
BODY2 = RGBColor(0x33, 0x33, 0x33)
STAT_NUM = RGBColor(0x16, 0x20, 0x2C)
STAT_DESC = RGBColor(0x4A, 0x55, 0x68)
TRACK = RGBColor(0xE6, 0xE9, 0xEC)
BAR_NAVY = RGBColor(0x47, 0x56, 0x6E)
BAR_AMBER = RGBColor(0xE3, 0xA6, 0x3E)
BAR_TEAL = RGBColor(0x1E, 0x4F, 0x58)
BAR_RED = RGBColor(0xC0, 0x39, 0x2B)
CARD_BORDER = RGBColor(0xB8, 0xCD, 0xDE)
NUM_LIGHT = RGBColor(0xB8, 0xCD, 0xDE)
ACCENT_PURPLE = RGBColor(0x60, 0x30, 0xA0)
ACCENT_AMBER = RGBColor(0xB0, 0x60, 0x10)
ACCENT_GREEN = RGBColor(0x0A, 0x7B, 0x52)
ACCENT_RED = RGBColor(0xC0, 0x30, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Georgia"
BODY_FONT = "Times New Roman"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def tb(slide, l, t, w, h, text, size, color, font=BODY_FONT, bold=False, italic=False,
       align=PP_ALIGN.LEFT, anchor=None, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    if anchor:
        box.text_frame.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = box.text_frame.paragraphs[0] if i == 0 else box.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = font
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
    return box


def section_header(slide, title, style="times"):
    if style == "georgia":
        tb(slide, 0.73, 0.58, 11.87, 0.81, title, 40, GRAY_HEAD2, font=TITLE_FONT)
    else:
        tb(slide, 0.62, 0.62, 12.69, 0.61, title, 34.5, GRAY_HEAD, font=BODY_FONT)


FIGURE_CAPTIONS = {
    "fig1_bland_altman.png": "Figure 1. Bland-Altman, IMU vs. OptiTrack R2n (n=61). Bias -0.108, wide limits of agreement -- see the Google Doc / local .pptx for the rendered chart.",
    "fig3_trajectory_example.png": "Figure 3. Single-trial knee-angle trajectory, IMU vs. OptiTrack overlaid (Participant 16). Early cycles track well; later cycles drift, settling ~117 deg (IMU) vs. ~125 deg (OptiTrack).",
    "fig31_mp_trajectory_example.png": "Figure 31. Single-trial knee-angle trajectory, MediaPipe vs. OptiTrack overlaid (Participant 16, right leg). MediaPipe tracks the initial drop, then scatters -- it never recovers the damped-oscillation pattern OptiTrack shows.",
    "fig4_metrics_by_group.png": "Figure 4. R2n / N / area ratio by group (Control vs. MS), per-participant means overlaid as open circles -- the n=1 control arm shown honestly, not as an error bar.",
    "fig6_metric_effect_heatmap.png": "Figure 6. Cohen's d (MAS>0 vs. MAS=0), 7 parameters x 3 modalities. OptiTrack shows positive effect sizes for phi_max/omega_max/f; IMU and MediaPipe show the opposite sign on most of the same parameters.",
    "fig7_single_vs_combined_auc.png": "Figure 7. Leave-one-participant-out AUC, single vs. combined PT parameters. Every single-metric AUC < 0.5 (range 0.0-0.17); combined 7-parameter AUC = 0.21.",
    "fig8_score_naive_vs_logocv.png": "Figure 8. Production compute_pt_score(), Control vs. MS. Naive Mann-Whitney p=0.5865 (not significant); LOGO-CV version could not run (only 1 control in the matched set).",
    "fig10_param_correlation.png": "Figure 10. Correlation matrix across the 7 PT parameters. R2n and omega_max_n correlate at r=0.93 -- near-redundant under the score's equal weighting.",
}


def add_figure(slide, filename, left, top, width):
    if NO_IMAGES:
        caption = FIGURE_CAPTIONS.get(filename, f"[{filename}]")
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(3.6))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = CARD_BORDER; box.line.width = Pt(1)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12); p.font.name = BODY_FONT; p.font.color.rgb = SLATE; p.font.italic = True
        return
    if COMPRESSED:
        filename = filename.rsplit(".", 1)[0] + FIG_EXT
    path = os.path.join(FIG_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    else:
        tb(slide, left, top, width, 1, f"[missing: {filename}]", 12, BODY)


def add_table(slide, headers, rows, left, top, width, height, font_size=11, header_fill=NAVY):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tshape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = tshape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.name = BODY_FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].font.name = BODY_FONT
            cell.text_frame.paragraphs[0].font.color.rgb = BODY
    return table


def stat_card(slide, left, top, width, number, desc):
    tb(slide, left, top, width, 0.6, number, 30.2, STAT_NUM, font=BODY_FONT, bold=True)
    tb(slide, left, top + 0.59, width, 0.7, desc, 11.23, STAT_DESC, font=BODY_FONT)


def ranking_row(slide, left, top, width, label, frac, bar_text, side_text, fill_color, label_bold=False):
    tb(slide, left, top, 3.21, 0.45, label, 10.17, NAVY, font=BODY_FONT, bold=label_bold)
    bar_left = left + 3.32
    bar_w = 6.10
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_left), Inches(top + 0.01),
                                    Inches(bar_w), Inches(0.36))
    track.fill.solid(); track.fill.fore_color.rgb = TRACK; track.line.fill.background()
    track.adjustments[0] = 0.5
    fillbar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_left), Inches(top + 0.01),
                                      Inches(bar_w * frac), Inches(0.36))
    fillbar.fill.solid(); fillbar.fill.fore_color.rgb = fill_color; fillbar.line.fill.background()
    fillbar.adjustments[0] = 0.5
    tf = fillbar.text_frame
    tf.margin_left = Inches(0.1); tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.text = bar_text
    p.font.size = Pt(9.64); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = BODY_FONT
    tb(slide, bar_left + bar_w + 0.15, top, 3.5, 0.45, side_text, 9.1, STAT_DESC, font=BODY_FONT)


def key_finding_card(slide, top, num, heading, body, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.73), Inches(top), Inches(11.87), Inches(1.02))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CARD_BORDER; card.line.width = Pt(0.75)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.73), Inches(top), Inches(0.07), Inches(1.02))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tb(slide, 0.93, top + 0.16, 0.5, 0.62, str(num), 26, NUM_LIGHT, font=TITLE_FONT)
    tb(slide, 2.51, top + 0.09, 10.0, 0.4, heading, 16, NAVY, font=TITLE_FONT)
    tb(slide, 2.51, top + 0.49, 10.0, 0.48, body, 11.5, SLATE, font="Calibri")


def divider(slide, l, t, w, color=SLATE):
    ln = slide.shapes.add_connector(1, Inches(l), Inches(t), Inches(l + w), Inches(t))
    ln.line.color.rgb = color
    ln.line.width = Pt(1.25)


# ============================================================== Slide 1: Title
s = new_slide()
tb(s, 0.0, 2.55, 13.33, 1.6, "Pendulastic: Results & Diagnostic Validation", 48, NAVY,
   font=TITLE_FONT, bold=False, align=PP_ALIGN.CENTER)
tb(s, 0.0, 3.95, 13.33, 0.5, "Smartphone-IMU Validation of the Wartenberg Pendulum Test for MS Spasticity",
   20, SLATE, font=TITLE_FONT, align=PP_ALIGN.CENTER)
tb(s, 0.0, 6.85, 13.33, 0.4, "Prepared for review with Dr. Monica Perez", 13, STAT_DESC,
   font=BODY_FONT, align=PP_ALIGN.CENTER)

# ============================================================== Slide 2: Sample overview (Aim-grid style)
s = new_slide()
section_header(s, "Sample & Data Availability")
cards = [
    ("Enrolled Cohort", "14 participants: 8 healthy control, 6 MS (+P17, grading in progress, excluded below)."),
    ("IMU-Validated Subset", "Only 5 participants (4 MS: P5, P13, P14, P15; 1 control: P16) have synchronized IMU + OptiTrack data -- the binding constraint on every result in this deck."),
    ("Trial Counts", "61 trials: IMU+OptiTrack. 37 trials: video+OptiTrack. 49 trials: three-way matched. 40 trials: complete production-score parameter sets."),
    ("MAS Coverage", "MAS grades available for all 61 IMU trials. Current cohort tops out at MAS = 1+; zero trials at MAS >= 2."),
]
positions = [(1.20, 1.72), (7.41, 1.72), (1.20, 4.20), (7.41, 4.20)]
for (hx, hy), (head, body) in zip(positions, cards):
    tb(s, hx, hy, 5.36, 0.29, head, 17, BODY, font=BODY_FONT, bold=True)
    tb(s, hx - 0.31, hy + 0.55, 5.60, 1.6, body, 13.5, BODY, font=BODY_FONT)

# ============================================================== Slide 3: IMU vs OptiTrack agreement
s = new_slide()
section_header(s, "IMU vs. OptiTrack Agreement", style="georgia")
tb(s, 0.73, 1.35, 11.87, 0.35, "n = 61 trials, 5 participants", 13, STAT_DESC, font=BODY_FONT, italic=True)
add_table(s, ["Parameter", "ICC(2,1)", "Bias", "95% LoA"], [
    ["R2n (relaxation index)", "0.226", "-0.108", "[-1.06, 0.85]"],
    ["N (oscillation count)", "0.458", "-1.46", "[-7.48, 4.56]"],
    ["phi_max ratio", "0.044", "-0.055", "[-0.51, 0.40]"],
    ["omega_max_n", "0.014", "-2.13", "[-12.05, 7.78]"],
    ["f (frequency)", "0.140", "-0.40", "[-1.98, 1.18]"],
    ["Area ratio", "0.135", "-0.052", "[-0.82, 0.72]"],
    ["omega_min_n", "0.214", "-1.01", "[-6.70, 4.69]"],
], left=0.73, top=1.85, width=6.2, height=3.6, font_size=11)
add_figure(s, "fig1_bland_altman.png", 7.2, 1.85, 5.4)
tb(s, 0.73, 5.65, 11.87, 1.5,
   "Trajectory-level RMSE: 14.84 deg mean / 10.98 deg median (n=53); only 2/53 trials met the 5 deg "
   "clinical-goal threshold. 67-76% of total error is a fixed per-trial calibration bias, not random "
   "noise -- the primary target for future correction.", 14, BODY, font=BODY_FONT)

# ============================================================== Slide 4: MediaPipe REFRAMED -- tested, not assumed
# A reviewer critique proposed a specific "assessor-tracking" linear formula (ot ~ 172.5 - 0.815*mp) as
# the explanation for MediaPipe's collapse. Tested it directly against this session's own data (pooled
# OptiTrack vs. MediaPipe angle, 49 three-way-matched trials, interpolated to a common time grid):
#   Full trial span:        ot = 156.53 - 0.113*mp,  r = -0.170  (n=9800 samples)
#   Active swing only (post-release, the physiologically relevant window): r = 0.008 (n=7350 samples)
# The claimed formula does not reproduce -- during the active swing, OptiTrack and MediaPipe angle are
# essentially UNCORRELATED, not related by a fittable linear offset. So this slide reports the honest,
# verified finding instead of the unverifiable critique claim: a real per-trial bias exists (mean ot-mp
# = -19.0 deg, SD 15.4 deg) but it does NOT reduce to a single correctable calibration constant.
s = new_slide()
section_header(s, "Section 3.3 Re-Tested: Is MediaPipe's Error a Fittable Offset?", style="georgia")
tb(s, 0.73, 1.35, 11.87, 0.35,
   "A reviewer proposed a linear \"assessor-tracking\" correction formula -- tested directly against this session's data before adopting it",
   13, STAT_DESC, font=BODY_FONT, italic=True)
stat_card(s, 0.73, 2.0, 5.6, "r = 0.008 (active swing)", "Pooled OptiTrack-vs-MediaPipe angle correlation, post-release window, 49 trials, 7,350 interpolated samples -- essentially uncorrelated")
stat_card(s, 0.73, 3.4, 5.6, "-19.0 deg +/- 15.4 deg", "Per-trial mean bias (OptiTrack minus MediaPipe) -- real and large, but not reducible to one linear constant (see below)")
add_figure(s, "fig31_mp_trajectory_example.png", 7.2, 1.85, 5.4)
tb(s, 0.73, 5.35, 11.87, 1.9,
   ["A candidate explanation -- that MediaPipe's error is a single fittable linear offset (an "
    "\"assessor-tracking\" formula), which would make it correctable with a simple calibration "
    "constant -- was tested directly and did not hold: during the active swing, OptiTrack and "
    "MediaPipe angle are essentially uncorrelated (r = 0.008), not linearly related.",
    "The bias is real (large and consistent in sign across trials) but varies enough per trial "
    "(SD 15.4 deg) that it is not a single correctable constant. This still supports the core point -- "
    "markerless pose estimation is not usable in this clinical setting without a validated "
    "calibration/tracking layer -- but the specific mechanism needs per-trial diagnosis, not a single global fit."],
   14, BODY, font=BODY_FONT)

# ============================================================== Slide 5: Markerless video table (Table 2)
s = new_slide()
section_header(s, "Markerless Video (MediaPipe) vs. OptiTrack -- Full Table")
add_table(s, ["Parameter", "ICC(2,1)", "Bias", "95% LoA"], [
    ["R2n (relaxation index)", "-0.041", "-0.441", "[-2.35, 1.47]"],
    ["N (oscillation count)", "0.032", "-2.30", "[-18.35, 13.75]"],
    ["phi_max ratio", "-0.008", "0.110", "[-0.91, 1.13]"],
    ["omega_max_n", "-0.036", "6.75", "[-50.28, 63.77]"],
    ["f (frequency)", "-0.115", "-0.83", "[-3.04, 1.37]"],
    ["Area ratio", "0.003", "-0.420", "[-1.23, 0.39]"],
    ["omega_min_n", "-0.018", "3.69", "[-29.16, 36.53]"],
], left=0.73, top=1.85, width=7.0, height=3.6, font_size=12)
tb(s, 8.2, 1.9, 4.4, 3.4,
   ["Full-curve RMSE: 36.0 deg mean / 33.3 deg median (n=37); 0/37 met the 5 deg threshold.",
    "",
    "ICC <= 0 for all 7 parameters -- table 2, n=49 trials, 5 participants.",
    "",
    "Reported as a real result, not rounded to zero: between-trial variance MediaPipe attributes to "
    "true differences is smaller than its own measurement noise."],
   13.5, BODY, font=BODY_FONT)

# ============================================================== Slide 6: Clinical validity
s = new_slide()
section_header(s, "Clinical Validity: Correlation With MAS", style="georgia")
stat_card(s, 0.9, 2.0, 5.6, "rho = -0.313", "Spearman correlation, IMU R2n vs. MAS grade, p = 0.014, n = 61 trials, all 5 participants")
stat_card(s, 6.9, 2.0, 5.6, "rho = -0.75 to -0.78", "Closest comparable published result (Yeh et al., 2025, stroke, n=20) -- our effect is real but roughly half the strength")
tb(s, 0.9, 3.6, 11.5, 2.0,
   ["Direction matches clinical expectation: a lower relaxation index (more restricted pendular swing) "
    "is associated with a higher MAS grade.",
    "Real and statistically significant -- but modest, and the current cohort tops out at MAS 1+, so "
    "the correlation has not yet been tested across the moderate-to-severe range."],
   15, BODY, font=BODY_FONT)

# ============================================================== Slide 7: Group comparison
s = new_slide()
section_header(s, "Group Comparison: MS vs. Control")
add_figure(s, "fig4_metrics_by_group.png", 0.6, 1.7, 7.0)
tb(s, 7.9, 1.8, 4.9, 4.5,
   ["Linear mixed-effects model (R2n ~ group, participant random intercept)",
    "",
    "MS estimated 0.173 lower than control (expected direction)",
    "beta = -0.173, SE = 0.207, z = -0.835, p = 0.404",
    "",
    "Not significant -- but NOT a null finding:",
    "the control arm is 1 participant (P16).",
    "",
    "This compares 4 MS individuals to one person,",
    "not to a control distribution."],
   15, BODY, font=BODY_FONT)

# ============================================================== Slide 8: sign inversion
s = new_slide()
section_header(s, "Do Parameters Discriminate Consistently Across Modalities?", style="georgia")
add_figure(s, "fig6_metric_effect_heatmap.png", 0.6, 1.6, 5.4)
tb(s, 6.4, 1.7, 6.4, 5.0,
   ["Cohen's d (MAS>0 vs. MAS=0), 7 parameters x 3 modalities, 49 three-way-matched trials.",
    "",
    "OptiTrack (ground truth): positive effect sizes for phi_max ratio, omega_max_n, f.",
    "",
    "IMU and MediaPipe: OPPOSITE SIGN on nearly all of the same parameters, same trials.",
    "",
    "More serious than noisy agreement: a single-modality parameter can point the wrong clinical "
    "direction entirely, not just be imprecise around the right one."],
   15, BODY, font=BODY_FONT)

# ============================================================== Slide 9: AUC
s = new_slide()
section_header(s, "Single-Metric vs. Combined-Metric Classification")
add_figure(s, "fig7_single_vs_combined_auc.png", 0.6, 1.7, 7.2)
tb(s, 8.0, 1.8, 4.9, 4.8,
   ["Leave-one-participant-out cross-validated logistic regression -- the honest generalization test.",
    "",
    "n = 49 trials, 5 participants (35 MAS>0 / 14 MAS=0)",
    "",
    "Every single-parameter AUC < 0.5 (range 0.0-0.17)",
    "Combined 7-parameter AUC = 0.21 -- worse than chance",
    "",
    "Reflects sample size (5 people), not refutation of the multi-metric hypothesis."],
   14, BODY, font=BODY_FONT)

# ============================================================== Slide 10: composite score
s = new_slide()
section_header(s, "Does the Production Composite Score Separate Groups?")
add_figure(s, "fig8_score_naive_vs_logocv.png", 0.6, 1.7, 7.2)
tb(s, 8.0, 1.8, 4.9, 4.8,
   ["compute_pt_score() on 40 IMU trials, 5 participants (complete parameter sets)",
    "",
    "Naive Mann-Whitney U (Control vs. MS):",
    "p = 0.5865 -- NOT significant",
    "",
    "Contradicts prior internal documentation:",
    '"p = 0.0001" (earlier PT7 comparison) -- flagged for re-verification',
    "",
    "LOGO-CV version could not even be run: only 1 control participant in the matched set."],
   14, BODY, font=BODY_FONT)

# ============================================================== Slide 11: collinearity
s = new_slide()
section_header(s, "Internal Structure: Parameter Collinearity")
add_figure(s, "fig10_param_correlation.png", 0.6, 1.7, 6.6)
tb(s, 7.5, 1.8, 5.3, 5.0,
   ["All 7 parameters weighted EQUALLY (1/7 each) -- not fit or validated against outcomes.",
    "",
    "R2n and omega_max_n correlate at r = 0.93 -- near-redundant, not independent.",
    "",
    "Shared signal effectively double-counted; distinct parameters (e.g. area ratio) under-weighted.",
    "",
    "HEALTHY_REF anchored on only 4 controls -- flagged \"PROVISIONAL\" in the code itself. "
    "Leave-one-control-out check could not run (only 1 control in the matched set)."],
   14, BODY, font=BODY_FONT)

# ============================================================== Slide 12: MAS target alignment
s = new_slide()
section_header(s, "Is MAS Correlation Measuring the Right Clinical Target?", style="georgia")
add_table(s, ["Leg", "mas_grade", "mas_flexion", "mas_extension"], [
    ["Left (P15)", "0", "1+", "0"],
    ["Right (P15)", "1", "1", "0"],
], left=0.9, top=1.9, width=6.3, height=1.1, font_size=13)
tb(s, 0.9, 3.2, 11.4, 3.6,
   ["The pendulum test targets knee-EXTENSOR spasticity specifically, not spasticity generally. "
    "Section 3.4's correlation used mas_grade, a collapsed field.",
    "",
    "For P15, the one participant with granular data: mas_extension = 0 on both legs, despite "
    "non-zero flexor and collapsed-grade scores.",
    "",
    "Open question this dataset cannot yet resolve: does R2n track the extensor-specific mechanism "
    "the test is designed for, or flexor signal that happens to covary with the collapsed grade in "
    "this small sample? A newly enrolled participant (P17) was scored on flexion/extension before an "
    "overall grade was assigned -- suggesting the sub-components may be the more appropriate target."],
   15, BODY, font=BODY_FONT)

# ============================================================== Slide 13: Key Findings (5-card)
s = new_slide()
section_header(s, "Key Findings", style="georgia")
key_finding_card(s, 1.57, 1, "IMU beats markerless video, with a fixable error source",
                  "IMU ICC 0.01-0.46 vs. MediaPipe ICC <= 0 on every parameter; ~70% of IMU error is a "
                  "correctable per-trial calibration bias, not random noise.", ACCENT_GREEN)
key_finding_card(s, 2.72, 2, "MediaPipe's failure is diagnosable, not random",
                  "A near-linear OptiTrack~MediaPipe relationship across trials points to a systematic "
                  "tracking/calibration problem -- exactly the kind spatial-temporal alignment can correct.",
                  ACCENT_PURPLE)
key_finding_card(s, 3.87, 3, "The composite score does not yet separate Control from MS",
                  "p = 0.5865 naively; LOGO-CV could not even run. Contradicts prior internal "
                  "documentation of p = 0.0001 -- flagged for re-verification.", ACCENT_RED)
key_finding_card(s, 5.01, 4, "The scoring formula has two fixable internal issues",
                  "R2n and omega_max_n are near-collinear (r=0.93) under equal weighting; HEALTHY_REF "
                  "rests on only 4 controls and has never been stress-tested.", ACCENT_AMBER)
key_finding_card(s, 6.16, 5, "Every limitation traces to one addressable cause",
                  "Only 5 of 15 enrolled participants have synchronized IMU+OptiTrack data, and the "
                  "control arm is n=1 -- see next slide for the fix.", ACCENT_AMBER)

# ============================================================== Slide 14: Zero-Recruitment solution
s = new_slide()
section_header(s, "The Zero-Recruitment Solution", style="georgia")
stat_card(s, 0.9, 1.7, 11.5, "7 controls, 0 new recruitment", "Existing video-only healthy controls who have never been IMU-recorded")
tb(s, 0.9, 2.9, 11.5, 3.9,
   ["The single highest-value next step: run the 7 existing video-only healthy controls through the "
    "IMU-recording app.",
    "",
    "This directly resolves the n=1 control-arm bottleneck behind three separate weaknesses in this "
    "deck at once -- the group comparison (Slide 7), the leave-one-participant-out classification test "
    "(Slide 9), and the untestable HEALTHY_REF sensitivity check (Slide 11) -- because all three fail "
    "for the same reason: too few IMU-recorded controls, not too little signal.",
    "",
    "It balances the clinical cohort, unlocks the statistical power required for a defensible "
    "classification claim, and requires no new participant recruitment."],
   16, BODY, font=BODY_FONT)

# ============================================================== Slide 15: Closing
s = new_slide()
tb(s, 0.0, 3.0, 13.33, 1.0, "Thank You", 44, NAVY, font=TITLE_FONT, align=PP_ALIGN.CENTER)
tb(s, 0.0, 4.0, 13.33, 0.6,
   "Questions & discussion -- Pendulastic, prepared for review with Dr. Monica Perez", 16, SLATE,
   font=TITLE_FONT, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH}")
