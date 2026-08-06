"""Generate Pendulastic Benchmark PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as In
from io import BytesIO
import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# ── Palette ────────────────────────────────────────────────────────────────
A   = RGBColor(0x00, 0x50, 0xA0)  # accent blue
G   = RGBColor(0x0A, 0x7B, 0x52)  # green (healthy / good)
GB  = RGBColor(0xDC, 0xF2, 0xE8)  # green bg
AM  = RGBColor(0xB0, 0x60, 0x10)  # amber (caution)
AB  = RGBColor(0xFD, 0xF0, 0xDC)  # amber bg
RD  = RGBColor(0xC0, 0x30, 0x30)  # red (bad)
RB  = RGBColor(0xFC, 0xEA, 0xEA)  # red bg
PV  = RGBColor(0x60, 0x30, 0xA0)  # purple (Pendulastic viewer)
PB  = RGBColor(0xED, 0xE0, 0xFA)  # purple bg
BL  = RGBColor(0x15, 0x65, 0xC0)  # blue
BB  = RGBColor(0xE0, 0xEA, 0xFA)  # blue bg
BG  = RGBColor(0xF0, 0xF5, 0xFA)  # slide background
T1  = RGBColor(0x0C, 0x1E, 0x34)  # dark text
T2  = RGBColor(0x3A, 0x58, 0x70)  # medium text
T3  = RGBColor(0x7A, 0x9C, 0xB8)  # muted text
BR  = RGBColor(0xB8, 0xCD, 0xDE)  # border grey-blue
PN  = RGBColor(0xE6, 0xEF, 0xF7)  # panel bg
WH  = RGBColor(0xFF, 0xFF, 0xFF)  # white

W = In(10)
H = In(5.625)

# ── Presentation setup ─────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

# ── Shape helpers ──────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def hline(slide, top, left=In(0.55), width=In(8.9), color=BR):
    sh = rect(slide, left, top, width, Pt(1), fill=color)
    return sh

# ── Text helpers ───────────────────────────────────────────────────────────

def tb(slide, text, l, t, w, h,
       fn='Calibri', fs=11, bold=False, italic=False,
       color=T1, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = wrap
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = fn; r.font.size = Pt(fs)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return box

def eyebrow(slide, text, top=In(0.28)):
    tb(slide, text.upper(), In(0.55), top, In(8.9), In(0.22),
       fn='Courier New', fs=8, bold=True, color=T3)

def title(slide, text, top=In(0.52), fs=30):
    tb(slide, text, In(0.55), top, In(8.9), In(0.8),
       fn='Georgia', fs=fs, color=T1)

def body_tf(slide, l, t, w, h):
    """Return a textframe-ready textbox with word wrap."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    # Clear default empty paragraph
    box.text_frame.paragraphs[0].clear()
    return box.text_frame

def add_p(tf, text, fn='Calibri', fs=10.5, bold=False, italic=False,
          color=T2, align=PP_ALIGN.LEFT, before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if before:
        p.space_before = Pt(before)
    r = p.add_run()
    r.text = text
    r.font.name = fn; r.font.size = Pt(fs)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color

def bullet_tf(slide, items, l, t, w, h,
              fn='Calibri', fs=10, color=T2, bullet_color=None, before=2):
    """Add a text box with bullet-point items."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = box.text_frame.paragraphs[0]
            first = False
        else:
            p = box.text_frame.add_paragraph()
        if before and not first:
            p.space_before = Pt(before)
        r = p.add_run()
        r.text = item
        r.font.name = fn; r.font.size = Pt(fs); r.font.color.rgb = color
    return box

def label_box(slide, label_text, body_text, l, t, w, fill_color, border_color):
    """Colored panel with a label header and body text."""
    rect(slide, l, t, w, In(0.68), fill=fill_color, line=border_color, lw=Pt(1))
    tb(slide, label_text.upper(), l + In(0.1), t + In(0.04), w - In(0.2), In(0.18),
       fn='Courier New', fs=7.5, bold=True, color=border_color)
    tf = body_tf(slide, l + In(0.1), t + In(0.22), w - In(0.2), In(0.42))
    add_p(tf, body_text, fs=9.5, color=T1)

# ── Step box for pipeline slides ──────────────────────────────────────────

def step_box(slide, num, num_color, num_bg, title_text, title_color,
             body_lines, l, t, w, h):
    """Numbered step box with title and body bullets."""
    rect(slide, l, t, w, h, fill=WH, line=BR, lw=Pt(0.75))
    # Number circle (approx — use small filled square)
    rect(slide, l + In(0.1), t + In(0.1), In(0.25), In(0.25), fill=num_bg)
    tb(slide, str(num), l + In(0.1), t + In(0.07), In(0.25), In(0.25),
       fn='Courier New', fs=8.5, bold=True, color=num_color,
       align=PP_ALIGN.CENTER)
    # Title
    tb(slide, title_text, l + In(0.42), t + In(0.08), w - In(0.55), In(0.22),
       fn='Calibri', fs=10, bold=True, color=title_color)
    # Body
    y_offset = In(0.32)
    box = slide.shapes.add_textbox(l + In(0.15), t + y_offset,
                                   w - In(0.25), h - y_offset - In(0.08))
    box.text_frame.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = box.text_frame.paragraphs[0]; first = False
        else:
            p = box.text_frame.add_paragraph()
            p.space_before = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.name = 'Calibri'
        r.font.size = Pt(9)
        r.font.color.rgb = T2

# ── Bar chart (shapes-based) ───────────────────────────────────────────────

def bar_chart(slide, data, l, t, chart_w, row_h=In(0.28),
              bar_max_w=In(4.5), label_w=In(1.7), max_val=100):
    """Horizontal bar chart using shapes. data = [(label, adj, exact, color, highlight)]"""
    for i, (label, adj, exact, bar_color, highlight) in enumerate(data):
        y = t + i * row_h
        # Label
        lc = PV if highlight else T2
        tb(slide, label, l, y + In(0.04), label_w - In(0.05), row_h - In(0.04),
           fn='Courier New', fs=8.5, bold=highlight, color=lc)
        # Adj bar (lighter)
        bw_adj = bar_max_w * min(adj, 100) / max_val
        if bw_adj > 0:
            rect(slide, l + label_w, y + In(0.04), bw_adj, row_h - In(0.1),
                 fill=bar_color if highlight else BL)
            # Make it lighter
            if not highlight:
                sh = rect(slide, l + label_w, y + In(0.04), bw_adj, row_h - In(0.1),
                          fill=BB)
        # Exact bar (darker, on top)
        bw_ex = bar_max_w * min(exact, 100) / max_val
        if bw_ex > 0:
            rect(slide, l + label_w, y + In(0.04), bw_ex, row_h - In(0.1),
                 fill=bar_color if highlight else G)
        # Value label
        val_text = f"{adj}% adj  {exact}% exact"
        tb(slide, val_text, l + label_w + max(bw_adj, In(0.05)) + In(0.05),
           y + In(0.05), In(2.5), row_h - In(0.08),
           fn='Courier New', fs=7.5, bold=highlight,
           color=PV if highlight else T3)
        # n=
        n_val = [m[0] for m in data].index(label)

# ── Slide builders ─────────────────────────────────────────────────────────

def s1_title(prs):
    slide = blank_slide(prs)
    tb(slide, 'PENDULASTIC RESEARCH  ·  ASU KINESIOLOGY',
       In(0.55), In(1.1), In(8), In(0.22), fn='Courier New', fs=7.5, color=T3)
    tb(slide, 'Pendulastic Benchmark',
       In(0.55), In(1.38), In(8.9), In(1.0), fn='Georgia', fs=42, color=T1)
    tb(slide, 'Pendulastic Viewer + HPE Models vs OptiTrack 3D Ground Truth',
       In(0.55), In(2.45), In(8), In(0.5), fn='Georgia', fs=16, color=T2)
    hline(slide, In(3.0), In(0.55), In(5.5))
    stats = [('11', 'systems compared'), ('9', 'participants'),
             ('241', 'HPE comparisons'), ('34', 'viewer sessions matched')]
    bw = In(1.65)
    for i, (v, lbl) in enumerate(stats):
        x = In(0.55) + i * bw
        rect(slide, x, In(3.15), bw - In(0.04), In(0.9), fill=WH, line=BR)
        tb(slide, v, x + In(0.12), In(3.2), bw - In(0.25), In(0.5),
           fn='Georgia', fs=30, color=A)
        tb(slide, lbl, x + In(0.12), In(3.72), bw - In(0.25), In(0.25),
           fn='Courier New', fs=7.5, color=T3)
    tb(slide, 'Identifying the most accurate markerless alternative to 3D motion capture '
       'for spasticity grading via the Wartenberg Pendulum Test',
       In(0.55), In(4.5), In(8.9), In(0.4),
       fn='Calibri', fs=9.5, italic=True, color=T3)


def s2_study_design(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'Study Design')
    title(slide, 'Participants & Protocol', top=In(0.50))
    hline(slide, In(1.08))

    # Left column — participant table
    col_l = In(0.55); col_w = In(4.3)
    tb(slide, 'PARTICIPANT SUMMARY', col_l, In(1.18), col_w, In(0.2),
       fn='Courier New', fs=7.5, bold=True, color=T3)

    headers = ['ID', 'Group', 'OptiTrack MAS', 'Viewer Sessions']
    col_xs  = [col_l, col_l+In(0.45), col_l+In(1.02), col_l+In(2.55)]
    col_ws  = [In(0.42), In(0.55), In(1.5), In(1.3)]
    header_y = In(1.38)
    for hdr, cx, cw in zip(headers, col_xs, col_ws):
        tb(slide, hdr, cx, header_y, cw, In(0.2),
           fn='Courier New', fs=7.5, bold=True, color=T3)
    hline(slide, In(1.58), col_l, col_w)

    rows = [
        ('P1',  'Control', 'MAS 0',        '—'),
        ('P2',  'Control', 'MAS 0–3',       '1 (left)'),
        ('P4',  'MS',      'MAS 0–2',       '5'),
        ('P5',  'MS',      'MAS 0 (all)*',  '11'),
        ('P6',  'Control', 'MAS 1–4 †',     '5'),
        ('P7',  'Control', 'MAS 0–1+',      '1'),
        ('P8',  'Control', 'MAS 0–1',       '7'),
        ('P9',  'Control', 'MAS 0–2',       '—'),
        ('P10', 'Control', 'MAS 0',         '5'),
    ]
    for ri, (pid, grp, mas, vs) in enumerate(rows):
        ry = In(1.62) + ri * In(0.34)
        is_ms = grp == 'MS'
        id_color  = AM if is_ms else T1
        mas_color = AM if is_ms else G
        vals = [pid, grp, mas, vs]
        colors = [id_color, T2, mas_color, T2]
        for val, cx, cw, col in zip(vals, col_xs, col_ws, colors):
            tb(slide, val, cx, ry, cw, In(0.3),
               fn='Courier New', fs=8.5, color=col)
        hline(slide, ry + In(0.3), col_l, col_w, color=PN)

    tb(slide, '* P5 paradox: MS patient scores MAS 0 on 3D kinematics\n'
              '† P6 has elevated muscle tone despite healthy classification\n'
              'P0 excluded — calibration reference only',
       col_l, In(4.85), col_w, In(0.55), fn='Calibri', fs=8, italic=True, color=T3)

    # Right column — protocol
    col2_l = In(5.1); col2_w = In(4.35)
    tb(slide, 'MEASUREMENT PROTOCOL', col2_l, In(1.18), col2_w, In(0.2),
       fn='Courier New', fs=7.5, bold=True, color=T3)

    protocols = [
        ('OptiTrack (Gold Standard)', A,
         '3D motion capture — quaternion knee angle from rigid body markers on thigh & shank. '
         '53 valid trials. Sub-millimeter accuracy.'),
        ('Pendulastic Viewer', PV,
         'Custom app — webcam video → 4-parameter Popovic pendulum score → estimated MAS grade. '
         '34 sessions matched to OptiTrack.'),
        ('HPE Models (10 families)', BL,
         'Standard webcam → neural network keypoint detection → knee angle → PT score → MAS. '
         '241 trial-model comparisons.'),
        ('MAS Scoring Scale', T2,
         'Popovic 4-parameter PT score → estimated MAS grade: 0 · 1 · 1+ · 2 · 3 · 4. '
         'Exact (same grade) and Adjacent (±1 step) agreement both reported.'),
    ]
    for pi, (lbl, lbl_color, desc) in enumerate(protocols):
        py = In(1.38) + pi * In(0.84)
        rect(slide, col2_l, py, col2_w, In(0.78), fill=WH, line=BR)
        tb(slide, lbl, col2_l + In(0.12), py + In(0.07), col2_w - In(0.2), In(0.2),
           fn='Courier New', fs=8.5, bold=True, color=lbl_color)
        tb(slide, desc, col2_l + In(0.12), py + In(0.3), col2_w - In(0.2), In(0.45),
           fn='Calibri', fs=9, color=T2)


def s3_optitrack(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'Technology — Ground Truth')
    title(slide, 'How OptiTrack Works', top=In(0.50))
    hline(slide, In(1.08))
    tb(slide, 'OptiTrack is a professional 3D motion capture system — the same technology used in movie VFX, sports science labs, and orthopedic clinics worldwide.',
       In(0.55), In(1.12), In(8.9), In(0.3), fn='Calibri', fs=9.5, italic=True, color=T3)

    steps = [
        ('1', 'Infrared Cameras',
         ['8-24 high-speed cameras mounted around the room',
          'Each emits near-infrared (NIR) light — invisible to human eyes (like a TV remote)',
          'Cameras capture at 120-360 frames per second',
          'Analogy: security cameras that only see reflective markers']),
        ('2', 'Retroreflective Markers',
         ['Small spheres (~12mm) coated in glass microspheres',
          'Attached to body at anatomical landmarks (thigh, shin)',
          'Reflect NIR light back to cameras — appear as bright dots',
          'Grouped into "rigid bodies": 3+ markers on one rigid plate per body segment']),
        ('3', '3D Triangulation',
         ['Each camera sees a marker from a different 2D angle',
          'Motive software cross-references all camera views',
          'Calculates exact 3D position in millimeters',
          'Same math as GPS — but with light rays. Accuracy: sub-millimeter.']),
        ('4', 'Quaternions & Rotation',
         ['Each rigid body rotation stored as quaternion [w, x, y, z]',
          '4 numbers that describe any 3D orientation without ambiguity',
          'Unlike pitch/roll/yaw, quaternions never "gimbal lock"',
          'Analogy: a compass bearing for 3D space with no blind spots']),
        ('5', 'Knee Angle & MAS',
         ['Thigh rigid body + Shank rigid body → angle between them = knee flex/ext angle',
          'Continuous time-series at 120-360 Hz with sub-mm accuracy',
          'Time-series → Popovic 4-param algorithm → PT score → MAS grade',
          'Software: NaturalPoint Motive']),
    ]

    n_steps = len(steps)
    sw = (In(8.9) - In(0.05) * (n_steps - 1)) / n_steps
    step_top = In(1.48)
    step_h   = In(2.88)

    for i, (num, ttl, bullets) in enumerate(steps):
        sx = In(0.55) + i * (sw + In(0.05))
        rect(slide, sx, step_top, sw, step_h, fill=WH, line=BR)
        # Step number badge
        rect(slide, sx + In(0.08), step_top + In(0.08), In(0.26), In(0.22), fill=BB)
        tb(slide, num, sx + In(0.08), step_top + In(0.06), In(0.26), In(0.24),
           fn='Courier New', fs=8.5, bold=True, color=A, align=PP_ALIGN.CENTER)
        # Title
        tb(slide, ttl, sx + In(0.08), step_top + In(0.33), sw - In(0.16), In(0.36),
           fn='Calibri', fs=9.5, bold=True, color=T1)
        hline(slide, step_top + In(0.72), sx + In(0.08), sw - In(0.16), color=PN)
        # Bullets
        box = slide.shapes.add_textbox(sx + In(0.08), step_top + In(0.78),
                                       sw - In(0.16), step_h - In(0.88))
        box.text_frame.word_wrap = True
        first = True
        for b in bullets:
            if first:
                p = box.text_frame.paragraphs[0]; first = False
            else:
                p = box.text_frame.add_paragraph()
                p.space_before = Pt(4)
            r = p.add_run()
            r.text = '• ' + b
            r.font.name = 'Calibri'; r.font.size = Pt(8.5); r.font.color.rgb = T2

    # Footer note
    rect(slide, In(0.55), In(4.44), In(8.9), In(0.72), fill=PN, line=BR)
    tb(slide, 'WHY IT\'S THE GOLD STANDARD', In(0.7), In(4.5), In(8.6), In(0.2),
       fn='Courier New', fs=7.5, bold=True, color=A)
    tb(slide, 'Direct 3D measurement — cannot be fooled by leg rotation toward/away from the camera (foreshortening). '
              'Used in biomechanics research, gait labs, sports science, and VFX worldwide. '
              'All HPE and Viewer scores in this study are benchmarked against OptiTrack as the definitive ground truth.',
       In(0.7), In(4.68), In(8.6), In(0.45), fn='Calibri', fs=9, color=T2)


def s4_hpe_tracking(prs):
    """Deep-dive: how HPE models estimate knee position + angle calculation."""
    slide = blank_slide(prs)
    eyebrow(slide, 'Technology — HPE Models')
    title(slide, 'How HPE Models Estimate Knee Position & Angle', top=In(0.50))
    hline(slide, In(1.08))

    # ── Left column: Neural Network pipeline ──────────────────────────────
    lx = In(0.55); lw = In(4.35)

    rect(slide, lx, In(1.18), lw, In(0.3), fill=BB)
    tb(slide, 'STEP 1-3  •  VIDEO → NEURAL NETWORK → KEYPOINTS',
       lx + In(0.1), In(1.22), lw - In(0.2), In(0.25),
       fn='Courier New', fs=8, bold=True, color=BL)

    nn_steps = [
        ('Input Frame',
         'A standard webcam or phone camera captures the test at ~30fps. '
         'Each frame is just an ordinary digital image — no special hardware.'),
        ('Convolutional Neural Network (CNN)',
         'A deep learning model trained on millions of labeled human images. '
         'It processes the image through dozens of layers:\n'
         '  Early layers detect edges & colors\n'
         '  Mid layers detect limb shapes & body regions\n'
         '  Final layers output a heatmap per joint\n'
         'Each heatmap is a "probability cloud" showing where a joint likely sits in the frame. '
         'The peak of the heatmap = predicted (x, y) pixel location.'),
        ('Keypoint Extraction',
         'Three keypoints define the knee angle: Hip, Knee, Ankle.\n'
         'These are extracted as (x, y) coordinates in pixels for every frame.\n'
         'Confidence scores filter unreliable detections (threshold: 0.3 – 0.5).\n'
         'Each model family uses a different architecture:\n'
         '  MediaPipe (Google) — BlazePose: 33 keypoints, optimized for mobile\n'
         '  RTMPose (Alibaba) — Real-Time Multi-person, COCO-trained\n'
         '  HRNet — High-Resolution Net: maintains spatial detail throughout\n'
         '  YOLO-Pose — Detection + pose in a single pass\n'
         '  DETRPose — Transformer-based detection\n'
         '  Pose2Sim — Multi-view 3D reconstruction from 2D views'),
    ]
    y = In(1.52)
    for ttl, desc in nn_steps:
        rect(slide, lx, y, lw, In(0.22), fill=PN)
        tb(slide, ttl, lx + In(0.1), y + In(0.03), lw - In(0.2), In(0.18),
           fn='Calibri', fs=9.5, bold=True, color=T1)
        # Body text — count lines to size box
        lines = desc.count('\n') + 1
        bh = In(0.18 * lines + 0.22)
        tb(slide, desc, lx + In(0.15), y + In(0.24), lw - In(0.25), bh,
           fn='Calibri', fs=8.5, color=T2)
        y += In(0.24) + bh + In(0.06)

    # ── Right column: Angle calculation ───────────────────────────────────
    rx = In(5.05); rw = In(4.4)

    rect(slide, rx, In(1.18), rw, In(0.3), fill=GB)
    tb(slide, 'STEP 4  •  HOW THE KNEE ANGLE IS CALCULATED',
       rx + In(0.1), In(1.22), rw - In(0.2), In(0.25),
       fn='Courier New', fs=8, bold=True, color=G)

    # Diagram — three points forming the angle
    # We use text to approximate a diagram
    diagram_box = rect(slide, rx, In(1.52), rw, In(1.55), fill=WH, line=BR)

    tb(slide, 'Hip (H)\n          ●\n          |\n          |  v1 = H - K\n          |\nKnee (K)●━━━━━━━━● Ankle (A)\n          v2 = A - K\n\n'
              'Angle θ = arccos( v1·v2 / |v1|·|v2| )\n'
              '        = arctan2( |v1 × v2|, v1·v2 )',
       rx + In(0.15), In(1.56), rw - In(0.3), In(1.45),
       fn='Courier New', fs=9, color=T1)

    calc_steps = [
        ('Vector construction',
         'v1 = Hip − Knee  (points from knee up to hip)\n'
         'v2 = Ankle − Knee  (points from knee down to ankle)'),
        ('Dot product formula',
         'cos θ = (v1 · v2) / (|v1| × |v2|)\n'
         'θ = arccos of the result → knee angle in degrees\n'
         'θ = 180° means fully straight leg. θ < 180° means bent.'),
        ('Frame-by-frame tracking',
         'This calculation runs on every video frame.\n'
         'Result: a continuous angle time-series (degrees vs. time)\n'
         'showing the leg swinging and oscillating after release.'),
        ('Into the Popovic algorithm',
         'The angle time-series is analysed for 4 biomechanical parameters:\n'
         '  R2n · N · phi_max_ratio · omega_max_n\n'
         'These parameters → PT score → estimated MAS grade.'),
    ]
    y2 = In(3.12)
    for ttl, desc in calc_steps:
        rect(slide, rx, y2, rw, In(0.2), fill=GB)
        tb(slide, ttl, rx + In(0.1), y2 + In(0.03), rw - In(0.2), In(0.17),
           fn='Calibri', fs=9.5, bold=True, color=G)
        lines = desc.count('\n') + 1
        bh = In(0.165 * lines + 0.15)
        tb(slide, desc, rx + In(0.15), y2 + In(0.22), rw - In(0.25), bh,
           fn='Calibri', fs=8.5, color=T2)
        y2 += In(0.22) + bh + In(0.05)

    # Footer — the 2D limitation
    rect(slide, In(0.55), In(4.95), In(8.9), In(0.55), fill=AB, line=AM, lw=Pt(1))
    tb(slide, 'THE 2D LIMITATION', In(0.7), In(4.99), In(4), In(0.2),
       fn='Courier New', fs=7.5, bold=True, color=AM)
    tb(slide, 'HPE models see only a 2D projection of the leg. If the leg rotates inward/outward (internal or external '
              'rotation), the 2D angle appears compressed — "foreshortening." OptiTrack measures the true 3D angle and '
              'is immune to this effect. This is the fundamental reason 2D-based methods are less accurate.',
       In(0.7), In(5.17), In(8.6), In(0.3), fn='Calibri', fs=8.5, color=AM)


def s5_pendulastic_algorithm(prs):
    """How the Pendulastic app works — pipeline and 4-parameter algorithm."""
    slide = blank_slide(prs)
    eyebrow(slide, 'Technology — Pendulastic App')
    title(slide, 'How Pendulastic Tracks & Scores the Pendulum Test', top=In(0.50))
    hline(slide, In(1.08))

    # Left: pipeline steps
    lx = In(0.55); lw = In(4.35)

    pipeline_steps = [
        (PV, PB, 'Video Input & HPE Keypoints',
         'Standard webcam or phone camera records the test. Pendulastic uses an HPE model internally to track the hip, knee, and ankle frame by frame — the same neural network approach described on the previous slide.'),
        (PV, PB, 'Knee Angle Time-Series',
         'The hip-knee-ankle angle is computed every frame using the vector dot product formula. The result is a waveform showing the leg swinging after it is released from full extension.'),
        (PV, PB, 'Swing Detection',
         'The algorithm detects when the leg crosses the resting (neutral) angle, counting each half-oscillation as one swing. It tracks until the leg settles within a small threshold of the resting position.'),
        (PV, PB, 'Popovic 4-Parameter Score',
         'Four biomechanical parameters are extracted from the waveform (see right column). These are combined into a single PT score that characterises how "freely" the leg swung.'),
        (PV, PB, 'Calibration → MAS Grade',
         'The PT score is compared against the healthy reference (P0). Scores near P0 map to MAS 0. Increasing deviation maps to MAS 1, 1+, 2, 3, 4. The MAS grade is displayed in real-time during the clinical exam.'),
    ]
    y = In(1.18)
    sh = (In(3.8)) / len(pipeline_steps)
    for i, (fc, bc, ttl, desc) in enumerate(pipeline_steps):
        rect(slide, lx, y, lw, sh - In(0.04), fill=bc, line=fc, lw=Pt(0.75))
        # Step number
        tb(slide, str(i + 1), lx + In(0.1), y + In(0.05), In(0.25), sh - In(0.1),
           fn='Courier New', fs=9, bold=True, color=fc)
        tb(slide, ttl, lx + In(0.38), y + In(0.05), lw - In(0.48), In(0.22),
           fn='Calibri', fs=9.5, bold=True, color=fc)
        tb(slide, desc, lx + In(0.38), y + In(0.27), lw - In(0.48), sh - In(0.35),
           fn='Calibri', fs=8.5, color=T1)
        y += sh

    # Real-time note
    rect(slide, lx, In(5.02), lw, In(0.45), fill=PB, line=PV, lw=Pt(1.5))
    tb(slide, 'Real-time output: MAS grade appears on screen during the examination — no post-processing. '
              'This is what makes Pendulastic a practical clinical tool rather than a research-only algorithm.',
       lx + In(0.12), In(5.06), lw - In(0.24), In(0.38),
       fn='Calibri', fs=9, italic=True, color=PV)

    # Right: 4-parameter deep dive
    rx = In(5.05); rw = In(4.4)
    rect(slide, rx, In(1.18), rw, In(0.3), fill=PB)
    tb(slide, 'THE 4-PARAMETER POPOVIC ALGORITHM',
       rx + In(0.1), In(1.22), rw - In(0.2), In(0.25),
       fn='Courier New', fs=8, bold=True, color=PV)

    params = [
        ('R2n  —  Oscillation Symmetry',
         'Measures how closely the leg\'s swing matches an ideal mathematical pendulum. '
         'A healthy leg swings symmetrically with equal amplitude on each side of neutral. '
         'Spasticity disrupts this symmetry — the "catch" pulls the leg differently on each swing. '
         'R2n close to 1.0 = healthy symmetric oscillation.'),
        ('N  —  Total Swing Count',
         'The number of full back-and-forth cycles before the leg settles to rest. '
         'A healthy leg swings freely many times (like a playground swing). '
         'A spastic leg has increased resistance — it settles in fewer swings. '
         'Lower N = more resistance = higher spasticity.'),
        ('phi_max_ratio  —  Peak Amplitude Ratio',
         'The ratio of the maximum swing amplitude to the initial drop angle. '
         'A healthy leg swings back up almost as high as it was dropped. '
         'A spastic leg barely swings back — the ratio is much smaller. '
         'phi_max_ratio close to 1.0 = healthy. Much less than 1.0 = spastic.'),
        ('omega_max_n  —  Normalised Peak Velocity',
         'The maximum angular velocity (how fast the leg swings at its fastest point), '
         'normalised by the expected velocity for a pendulum of that length. '
         'Spasticity slows the peak velocity disproportionately. '
         'Combined with N and phi_max_ratio, this gives a complete picture of swing dynamics.'),
    ]
    py = In(1.52)
    ph = In(0.84)
    for i, (param_name, param_desc) in enumerate(params):
        rect(slide, rx, py, rw, ph, fill=WH, line=BR)
        tb(slide, param_name, rx + In(0.12), py + In(0.07), rw - In(0.24), In(0.22),
           fn='Courier New', fs=9, bold=True, color=PV)
        tb(slide, param_desc, rx + In(0.12), py + In(0.3), rw - In(0.24), ph - In(0.36),
           fn='Calibri', fs=8.5, color=T2)
        py += ph + In(0.04)

    # MAS calibration note
    rect(slide, rx, In(4.94), rw, In(0.55), fill=PN, line=BR)
    tb(slide, 'CALIBRATION', rx + In(0.12), In(4.98), rw - In(0.24), In(0.2),
       fn='Courier New', fs=7.5, bold=True, color=T3)
    tb(slide, 'PT score is mapped to MAS using a lookup table calibrated against P0 '
              '(healthy reference participant). P0\'s score sets the MAS 0 baseline; '
              'increasing deviation maps to MAS 1 → 1+ → 2 → 3 → 4.',
       rx + In(0.12), In(5.16), rw - In(0.24), In(0.3),
       fn='Calibri', fs=8.5, color=T2)


def s_param_profile(prs):
    """MS vs Control parameter profiles — which parameters matter most for MS detection."""
    slide = blank_slide(prs)

    # ── Header (compact — title fs=17 to leave room for table) ────────────────
    eyebrow(slide, 'Parameter Analysis — MS vs Control')
    tb(slide, 'Ranking the 7 Pendulum Parameters: Which Best Identify MS?',
       In(0.55), In(0.44), In(9.0), In(0.52),
       fn='Georgia', fs=22, color=T1)
    tb(slide,
       'Mann-Whitney U tests  ·  n=13 control trials  ·  n=8 MS trials  ·  '
       'Cliff\'s delta: effect size (0-1, higher = stronger separation)',
       In(0.55), In(0.96), In(9.0), In(0.18),
       fn='Courier New', fs=7.5, color=T3)
    hline(slide, In(1.16))

    # ── Layout constants ───────────────────────────────────────────────────────
    HL = In(0.55)   # left edge
    RH = In(0.40)   # row height
    # Column start positions & widths
    #  [0] tier label  [1] param name   [2] direction  [3] bar  [4] delta  [5] p  [6] ctrl  [7] ms
    CX = [HL,       In(0.95), In(3.15), In(3.78), In(5.22), In(5.74), In(6.58), In(7.38)]
    CW = [In(0.38), In(2.16), In(0.58), In(1.38), In(0.48), In(0.80), In(0.76), In(2.0)]
    BAR_ORIGIN = CX[3]
    BAR_MAX    = CW[3] - In(0.06)

    # ── Column headers ─────────────────────────────────────────────────────────
    HDRS = ['', 'Parameter', 'MS vs ctrl', "Cliff's delta (effect size bar)", '|delta|', 'p-value', 'Control', 'MS']
    hy = In(1.20)
    rect(slide, HL, hy, In(8.9), In(0.22), fill=PN)
    for hdr, cx, cw in zip(HDRS, CX, CW):
        tb(slide, hdr, cx + In(0.03), hy + In(0.03), cw - In(0.04), In(0.17),
           fn='Courier New', fs=6.5, bold=True, color=T3)

    # ── Row data ───────────────────────────────────────────────────────────────
    # tier, name, direction, delta, p, ctrl_str, ms_str, note
    rows = [
        ('PRIMARY',       'phi_max  (A2/A0)',     'LOWER',  0.981, 0.0002, '0.79',       '0.13',
         'How high the leg rebounds on first return swing. Near-zero in MS = barely any swing-back. Most discriminating single parameter.'),
        ('PRIMARY',       'N  (swing count)',       'LOWER',  0.875, 0.001,  '7.0 swings', '2.3 swings',
         'Total full oscillations before leg settles. Controls swing freely 7+ times; MS patients stop in 2-3 swings.'),
        ('SECONDARY',     'R2n  (relaxation)',      'LOWER',  0.519, 0.055,  '1.01',       '0.96',
         'Symmetry of the first oscillation (A1 peak-to-peak / 1.6xA0). Small but consistent reduction in MS; borderline p.'),
        ('SECONDARY',     'omega_max  (peak vel.)', 'LOWER',  0.481, 0.076,  '5.69 /s',    '4.88 /s',
         'Max angular speed normalised to drop angle. MS leg moves slightly slower; supports phi_max and N but not significant alone.'),
        ('SUPPLEMENTARY', 'f  (frequency Hz)',      'SIMILAR',0.0,   None,   '~0.90 Hz',   '~0.85 Hz',
         'Natural oscillation rate. Depends on leg length; excluded from primary score (population-specific reference).'),
        ('SUPPLEMENTARY', 'area_ratio  (symmetry)', 'SIMILAR',0.0,   None,   '~0.09',      'varies',
         'Asymmetry of positive vs negative swing area. Useful quality flag for atypical recordings; high trial-to-trial variability.'),
        ('EXCLUDED',      'omega_min  (min vel.)',  'SIMILAR',0.0,   None,   '~0.001',     '~0.000',
         'Minimum angular speed during oscillation. Approaches zero for ALL subjects as swings decay — non-discriminating, excluded from scoring.'),
    ]

    tier_fg = {'PRIMARY': G,  'SECONDARY': AM, 'SUPPLEMENTARY': BL, 'EXCLUDED': T3}
    tier_bg = {'PRIMARY': GB, 'SECONDARY': AB, 'SUPPLEMENTARY': BB, 'EXCLUDED': PN}
    dir_col = {'LOWER': RD, 'HIGHER': AM, 'SIMILAR': T3}
    dir_sym = {'LOWER': 'lower  v', 'HIGHER': 'higher  ^', 'SIMILAR': 'similar'}

    ry = hy + In(0.24)
    for tier, name, direction, delta, p, ctrl_str, ms_str, note in rows:
        fc  = tier_fg[tier]
        bc  = tier_bg[tier]
        is_primary   = tier == 'PRIMARY'
        is_excluded  = tier == 'EXCLUDED'
        row_fill     = bc if is_primary else (PN if is_excluded else WH)
        txt_color    = fc if is_primary else (T3 if is_excluded else T2)

        # Row background
        rect(slide, HL, ry, In(8.9), RH - Pt(1.5), fill=row_fill, line=BR, lw=Pt(0.3))
        # Left accent stripe for primary rows
        if is_primary:
            rect(slide, HL, ry, Pt(4), RH - Pt(1.5), fill=fc)

        # [0] Tier label (small vertical badge)
        rect(slide, CX[0] + In(0.02), ry + In(0.06), CW[0] - In(0.04), RH - In(0.18), fill=bc)
        tier_abbr = tier[:3]
        tb(slide, tier_abbr, CX[0] + In(0.02), ry + In(0.07), CW[0] - In(0.04), In(0.18),
           fn='Courier New', fs=6, bold=True, color=fc, align=PP_ALIGN.CENTER)

        # [1] Parameter name
        tb(slide, name, CX[1] + In(0.04), ry + In(0.09), CW[1] - In(0.06), In(0.22),
           fn='Calibri', fs=9.5, bold=is_primary, color=txt_color)

        # [2] MS direction
        tb(slide, dir_sym[direction], CX[2] + In(0.04), ry + In(0.09), CW[2] - In(0.06), In(0.22),
           fn='Courier New', fs=8, bold=is_primary, color=dir_col[direction])

        # [3] Effect size bar
        if delta > 0:
            bw = BAR_MAX * delta
            bar_fill = G if is_primary else (AM if tier == 'SECONDARY' else BL)
            rect(slide, BAR_ORIGIN, ry + In(0.09), bw, RH - In(0.22), fill=bar_fill)
        else:
            # placeholder tick
            rect(slide, BAR_ORIGIN, ry + In(0.16), In(0.05), RH - In(0.32), fill=BR)

        # [4] Delta value
        delta_str = f'{delta:.3f}' if delta > 0 else '—'
        tb(slide, delta_str, CX[4] + In(0.03), ry + In(0.09), CW[4] - In(0.04), In(0.22),
           fn='Courier New', fs=8, bold=is_primary, color=fc if is_primary else txt_color)

        # [5] p-value
        if p is not None:
            if p < 0.001:   p_str, p_c = 'p<0.001', G
            elif p < 0.01:  p_str, p_c = f'p={p:.3f}', G
            elif p < 0.05:  p_str, p_c = f'p={p:.3f}', G
            elif p < 0.1:   p_str, p_c = f'p={p:.3f}', AM
            else:           p_str, p_c = f'p={p:.3f}', T3
        else:
            p_str, p_c = 'not tested', T3
        tb(slide, p_str, CX[5] + In(0.03), ry + In(0.09), CW[5] - In(0.04), In(0.22),
           fn='Courier New', fs=8, bold=(p is not None and p < 0.05), color=p_c)

        # [6] Control median
        tb(slide, ctrl_str, CX[6] + In(0.03), ry + In(0.09), CW[6] - In(0.04), In(0.22),
           fn='Courier New', fs=8, color=G)

        # [7] MS median + note (note wraps below in same column space)
        ms_c = RD if direction == 'LOWER' else (AM if direction == 'HIGHER' else T3)
        tb(slide, ms_str, CX[7] + In(0.04), ry + In(0.09), CW[7] - In(0.06), In(0.22),
           fn='Courier New', fs=8, bold=is_primary, color=ms_c)

        ry += RH

    # ── Right-side note column: placed as a separate overlapping column ────────
    # Re-draw notes in a dedicated panel to the right of the table
    # (We re-iterate the rows at the same y positions)
    NOTE_L = In(7.38); NOTE_W = In(2.05)
    NOTE_START = hy + In(0.24)
    for i, (tier, name, direction, delta, p, ctrl_str, ms_str, note) in enumerate(rows):
        ny = NOTE_START + i * RH
        is_excluded = tier == 'EXCLUDED'
        tb(slide, note, NOTE_L + In(0.05), ny + In(0.04), NOTE_W - In(0.08), RH - In(0.08),
           fn='Calibri', fs=7.5, color=T3 if is_excluded else T2)

    # ── Footer: legend + clinical thresholds ───────────────────────────────────
    FOOT_TOP = ry + In(0.08)
    FOOT_H   = In(5.45) - FOOT_TOP

    # Left: legend
    rect(slide, HL, FOOT_TOP, In(4.3), FOOT_H, fill=WH, line=BR)
    tb(slide, 'KEY', HL + In(0.12), FOOT_TOP + In(0.07), In(0.5), In(0.16),
       fn='Courier New', fs=7, bold=True, color=T3)
    legend = [
        (G,  GB, 'PRIMARY',       'Significant (p<0.01), large effect — check these first'),
        (AM, AB, 'SECONDARY',     'Trending (p<0.1) — supporting, not conclusive alone'),
        (BL, BB, 'SUPPLEMENTARY', 'Tracked but less discriminating for MS'),
        (T3, PN, 'EXCLUDED',      'Near-zero for all subjects — removed from scoring'),
    ]
    ly = FOOT_TOP + In(0.07)
    for fc, bc, lbl, desc in legend:
        rect(slide, HL + In(0.12), ly + In(0.18), In(0.44), In(0.16), fill=bc, line=fc, lw=Pt(0.5))
        tb(slide, lbl[:3], HL + In(0.12), ly + In(0.19), In(0.44), In(0.14),
           fn='Courier New', fs=6.5, bold=True, color=fc, align=PP_ALIGN.CENTER)
        tb(slide, desc, HL + In(0.6), ly + In(0.19), In(3.6), In(0.16),
           fn='Calibri', fs=8, color=T2)
        ly += In(0.22)

    # Right: clinical thresholds
    rect(slide, In(4.5), FOOT_TOP, In(5.0), FOOT_H, fill=GB, line=G, lw=Pt(1.5))
    tb(slide, 'CLINICAL DECISION GUIDE  —  START WITH phi_max AND N',
       In(4.62), FOOT_TOP + In(0.07), In(4.78), In(0.18),
       fn='Courier New', fs=7.5, bold=True, color=G)
    hline(slide, FOOT_TOP + In(0.27), In(4.62), In(4.78), color=G)
    thresholds = [
        ('phi_max < 0.30',   'strong spasticity signal',       'controls average 0.79'),
        ('N  <  4 swings',   'strong spasticity signal',       'controls average 7.0'),
        ('PT score > 0.28',  'estimated MAS 1 or above',       '(composite of all 4 primary params)'),
        ('R2n  <  0.85',     'supportive evidence',            'use alongside phi_max + N'),
    ]
    ty = FOOT_TOP + In(0.32)
    for trigger, action, sub in thresholds:
        tb(slide, trigger, In(4.62), ty, In(1.35), In(0.18),
           fn='Courier New', fs=8.5, bold=True, color=T1)
        tb(slide, f'=> {action}', In(6.0), ty, In(1.75), In(0.18),
           fn='Courier New', fs=8.5, color=G)
        tb(slide, sub, In(7.78), ty, In(1.6), In(0.18),
           fn='Courier New', fs=7.5, color=T3)
        ty += In(0.22)


def s6_viewer_results(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'Pendulastic App Performance')
    title(slide, 'Pendulastic Viewer vs OptiTrack Ground Truth', top=In(0.50))
    hline(slide, In(1.08))

    # Hero number
    rect(slide, In(0.55), In(1.18), In(4.0), In(1.8), fill=PB, line=PV, lw=Pt(2))
    tb(slide, '85.3%', In(0.7), In(1.25), In(3.8), In(1.0),
       fn='Georgia', fs=56, bold=False, color=PV)
    tb(slide, 'Adjacent MAS Agreement (within ±1 grade)',
       In(0.7), In(2.25), In(3.8), In(0.3), fn='Calibri', fs=11, color=T2)
    tb(slide, '52.9% exact match  ·  34 matched sessions  ·  5 non-adjacent errors',
       In(0.7), In(2.56), In(3.8), In(0.3), fn='Courier New', fs=8.5, color=T3)

    # Stats row
    stats = [('52.9%', 'Exact MAS match'), ('0.394', 'PT score r'), ('0.100', 'Mean abs PT error'), ('#2', 'Overall rank')]
    sw = In(1.1)
    for i, (v, lbl) in enumerate(stats):
        sx = In(0.55) + i * (sw + In(0.05))
        tb(slide, v, sx, In(3.08), sw, In(0.38), fn='Georgia', fs=22, color=T1)
        tb(slide, lbl, sx, In(3.48), sw, In(0.25), fn='Courier New', fs=7.5, color=T3)

    # By population
    rect(slide, In(0.55), In(3.82), In(1.92), In(0.72), fill=GB, line=G, lw=Pt(1))
    tb(slide, '57.9%', In(0.65), In(3.88), In(1.7), In(0.38), fn='Georgia', fs=24, color=G)
    tb(slide, 'Healthy → MAS 0\n(11/19 trials)', In(0.65), In(4.27), In(1.7), In(0.25),
       fn='Calibri', fs=8.5, color=T2)

    rect(slide, In(2.57), In(3.82), In(1.98), In(0.72), fill=AB, line=AM, lw=Pt(1))
    tb(slide, '53.3%', In(2.67), In(3.88), In(1.8), In(0.38), fn='Georgia', fs=24, color=AM)
    tb(slide, 'MS → MAS ≥1\n(8/15 MS trials)', In(2.67), In(4.27), In(1.8), In(0.25),
       fn='Calibri', fs=8.5, color=T2)

    # Comparison panel
    rect(slide, In(0.55), In(4.62), In(3.95), In(0.88), fill=PN, line=BR)
    tb(slide, 'VS TOP HPE MODELS', In(0.68), In(4.67), In(3.7), In(0.2),
       fn='Courier New', fs=7.5, bold=True, color=T3)
    rows_cmp = [
        ('Pendulastic Viewer  ★', '85.3% adj', '57.9% hlth→0', '53.3% MS→1+', PV),
        ('MediaPipe Full',        '84.4% adj', '37.5% hlth→0', '87.5% MS→1+', T2),
        ('MediaPipe Lite',        '83.3% adj', '16.7% hlth→0', '91.7% MS→1+', T2),
    ]
    for ri, (name, adj, h0, ms, col) in enumerate(rows_cmp):
        ry = In(4.87) + ri * In(0.2)
        tb(slide, name,  In(0.68), ry, In(1.5), In(0.19), fn='Courier New', fs=8.5, bold=(col==PV), color=col)
        tb(slide, adj,   In(2.2),  ry, In(0.7), In(0.19), fn='Courier New', fs=8.5, color=col)
        tb(slide, h0,    In(2.95), ry, In(0.85),In(0.19), fn='Courier New', fs=8.5, color=col)
        tb(slide, ms,    In(3.84), ry, In(0.62),In(0.19), fn='Courier New', fs=8.5, color=col)

    # Right: PT correlation note
    rect(slide, In(4.65), In(1.18), In(4.8), In(4.32), fill=WH, line=BR)
    tb(slide, 'WHY THE VIEWER\'S PT CORRELATION IS STRONGER',
       In(4.78), In(1.25), In(4.52), In(0.3),
       fn='Courier New', fs=8, bold=True, color=PV)
    notes = [
        'The Pendulastic Viewer achieves r = 0.394 PT score correlation with OptiTrack — stronger than MediaPipe Full (r = 0.156) and most other HPE models.',
        'This is because the 4-parameter Popovic algorithm directly captures the biomechanical features of the pendulum swing — symmetry, count, amplitude, and velocity — rather than just looking at raw angle values.',
        'HPE models compute a PT score from the same angle time-series, but 2D foreshortening errors add noise to every parameter, dragging the correlation down.',
        'The Pendulastic app\'s real advantage is in the algorithm design, not just the keypoint detection.',
        '',
        'HEALTHY CLASSIFICATION ADVANTAGE',
        'The Viewer correctly identifies 57.9% of healthy trials as MAS 0, compared to only 37.5% for MediaPipe Full and 16.7% for MediaPipe Lite.',
        'HPE models tend to over-grade spasticity in healthy participants — probably because they overcount swing cycles (N bias +1.5 to +4.0 cycles), which the Popovic algorithm interprets as reduced oscillation.',
        '',
        'THE MS DETECTION GAP',
        'Only 53.3% of MS trials scored MAS ≥1 by the Viewer, vs 87.5% for MediaPipe Full.',
        'P5 explains much of this: OptiTrack scores all P5 trials MAS 0 (no pendulum spasticity), while the Viewer and HPE models detect tone — likely reflecting real subclinical spasticity the 3D kinematics miss.',
    ]
    y_n = In(1.58)
    for note in notes:
        if note == '':
            y_n += In(0.1)
            continue
        is_header = note[0].isupper() and len(note) < 50 and note == note.upper() or note.endswith('ADVANTAGE') or note.endswith('GAP') or note.endswith('STRONGER')
        # Detect section headers
        is_hdr = note.startswith('HEALTHY') or note.startswith('THE MS') or note.startswith('WHY')
        clr = PV if is_hdr else T2
        bld = is_hdr
        fs  = 8 if is_hdr else 9
        box = slide.shapes.add_textbox(In(4.78), y_n, In(4.52), In(0.4))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = note
        r.font.name = 'Calibri'; r.font.size = Pt(fs)
        r.font.bold = bld; r.font.color.rgb = clr
        lines = max(1, len(note) // 70 + 1)
        y_n += In(0.155 * lines + 0.02)


def s7_rankings(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'All Systems')
    title(slide, 'Accuracy Rankings — Adjacent & Exact MAS Agreement', top=In(0.50))
    hline(slide, In(1.08))

    # Legend
    tb(slide, '★  Pendulastic Viewer (our app)       †  n < 15 (limited data)       ⚠  systematic bias       ✕  not viable',
       In(0.55), In(1.12), In(8.9), In(0.25),
       fn='Courier New', fs=8, color=T3)

    # Table header
    cols = ['System / Model', 'n', 'Adj %', 'Exact %', 'Hlth→0%', 'MS→1+%', 'r (PT)', 'MAE']
    col_xs = [In(0.55), In(3.05), In(3.45), In(4.1), In(4.78), In(5.52), In(6.26), In(7.0)]
    col_ws = [In(2.45), In(0.36), In(0.6),  In(0.64),In(0.68), In(0.68), In(0.68), In(0.68)]
    hy = In(1.38)
    hline(slide, hy + In(0.22), In(0.55), In(8.9))
    for hdr, cx, cw in zip(cols, col_xs, col_ws):
        tb(slide, hdr, cx, hy, cw, In(0.22),
           fn='Courier New', fs=7.5, bold=True, color=T3)

    rows = [
        ('pose2sim_lw †',      '13', '92.3', '53.8', '57.1', '83.3', '0.773', '0.099',  False, 'limited'),
        ('★ Pendulastic Viewer','34', '85.3', '52.9', '57.9', '53.3', '0.394', '0.100',  True,  'viewer'),
        ('MediaPipe Full',      '32', '84.4', '53.1', '37.5', '87.5', '0.156', '0.152',  False, None),
        ('MediaPipe Lite',      '42', '83.3', '40.5', '16.7', '91.7', '0.265', '0.138',  False, None),
        ('MediaPipe Heavy',     '41', '75.6', '36.6', '33.3', '81.8', '0.073', '0.163',  False, None),
        ('RTMPose-S',           '16', '68.8', '31.2', '50.0', '75.0', '0.395', '0.168',  False, None),
        ('RTMPose-M',           ' 9', '66.7', '22.2', '0.0',  '66.7', '0.151', '0.230',  False, None),
        ('HRNet-W48 ⚠',          '26', '65.4', '26.9', '33.3', '92.9', '0.117', '0.230',  False, 'caution'),
        ('DETRPose-S',          '29', '37.9', '10.3', '11.5','100.0', '0.160', '0.432',  False, None),
        ('YOLO-N ✕',             '32', ' 6.2', ' 3.1', ' 0.0','100.0', '0.061', '0.578',  False, 'fail'),
        ('RTMO-M †',             ' 1','100.0', ' 0.0',  '—',  '100.0',  '—',   '0.027',  False, 'limited'),
    ]
    ry = In(1.64)
    rh = In(0.32)
    for row in rows:
        name, n, adj, exact, h0, ms, r, mae, is_viewer, tag = row
        bg = PB if is_viewer else (AB if tag=='caution' else (RB if tag=='fail' else WH))
        nc = PV if is_viewer else (AM if tag=='caution' else (RD if tag=='fail' else (T3 if tag=='limited' else T2)))
        rect(slide, In(0.55), ry, In(8.9), rh - In(0.02), fill=bg, line=BR, lw=Pt(0.4))
        vals = [name, n, adj, exact, h0, ms, r, mae]
        for val, cx, cw in zip(vals, col_xs, col_ws):
            is_name = val == name
            tb(slide, str(val), cx + In(0.05), ry + In(0.06), cw - In(0.08), rh - In(0.1),
               fn='Courier New', fs=8.5, bold=(is_viewer and is_name),
               color=(PV if is_viewer else nc if is_name else T2))
        ry += rh

    # Bar chart visual (simplified — adj% bars using shapes)
    bar_top = ry + In(0.08)
    bar_l = In(3.45); bar_max = In(5.5)
    tb(slide, 'Adjacent agreement %', bar_l, bar_top - In(0.2), bar_max, In(0.18),
       fn='Courier New', fs=7.5, color=T3)
    for i, row in enumerate(rows):
        name, n, adj, *_ , is_viewer, tag = row
        try:
            adj_f = float(adj.strip())
        except:
            adj_f = 0
        bw = bar_max * adj_f / 100
        fill_c = PV if is_viewer else (RD if tag=='fail' else (AM if tag=='caution' else BL))
        by = bar_top + i * In(0.28)
        if bw > 0:
            rect(slide, bar_l, by + In(0.04), bw, In(0.2), fill=fill_c)
        tb(slide, f'{adj.strip()}%', bar_l + bw + In(0.04), by + In(0.05), In(0.6), In(0.2),
           fn='Courier New', fs=7.5, color=T3)


def s8_per_participant(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'Participant Analysis')
    title(slide, 'Pendulastic Viewer — Per-Participant Breakdown', top=In(0.50))
    hline(slide, In(1.08))

    cols = ['Participant', 'Group', 'OptiTrack MAS', 'Viewer Exact/Adj/Total', 'Key Observation']
    col_xs = [In(0.55), In(1.45), In(2.12), In(3.35), In(4.95)]
    col_ws = [In(0.85), In(0.62), In(1.18), In(1.55),  In(4.5)]
    hy = In(1.18)
    for hdr, cx, cw in zip(cols, col_xs, col_ws):
        tb(slide, hdr, cx, hy, cw, In(0.22), fn='Courier New', fs=7.5, bold=True, color=T3)
    hline(slide, In(1.42))

    rows = [
        ('P2',  'Healthy', 'MAS 0',        1, 1, 1,  G,  GB, 'Single trial — perfect agreement on healthy knee.'),
        ('P4',  'MS',      'MAS 0–1+',     1, 4, 5,  AM, AB, 'Spasticity detected in 4/5 trials (adj). 1 session scored MAS 0 vs OptiTrack 1+. Duplicate viewer sessions show within-session variability.'),
        ('P5',  'MS',      'MAS 0 (all)',  5, 7, 11, AM, AB, 'P5 PARADOX: OptiTrack shows MAS 0 every trial, yet Viewer scores MAS 1-3 on 6 sessions. May reflect real clinical tone that pendulum kinematics miss. 3 non-adjacent misses.'),
        ('P6',  'Healthy', 'MAS 1–4',      2, 4, 5,  BL, BB, 'Highest-tone healthy participant. 4/5 adjacent. One outlier: Viewer scored MAS 3 vs OptiTrack MAS 1 (right T2).'),
        ('P7',  'Healthy', 'MAS 0–1+',     0, 1, 1,  BL, BB, 'Single trial — Viewer scored MAS 1 vs OptiTrack MAS 0. Adjacent, borderline case.'),
        ('P8',  'Healthy', 'MAS 0–1',      5, 7, 7,  G,  GB, 'Best performing participant. All 7 trials adjacent, 5 exact. Cleanest pendulum data in the cohort.'),
        ('P10', 'Healthy', 'MAS 0',        4, 5, 5,  G,  GB, 'Near-perfect. Viewer scored MAS 1 on one left-leg trial vs OptiTrack MAS 0. Single adjacent miss.'),
    ]
    ry = In(1.48)
    for pid, grp, mas, e, a, total, fc, bc, note in rows:
        is_ms = grp == 'MS'
        rh = In(0.58) if len(note) > 100 else In(0.48)
        rect(slide, In(0.55), ry, In(8.9), rh, fill=bc if is_ms else WH, line=BR, lw=Pt(0.4))
        if is_ms:
            rect(slide, In(0.55), ry, Pt(3), rh, fill=AM)
        tb(slide, pid,  col_xs[0], ry+In(0.1), col_ws[0], rh-In(0.15), fn='Courier New', fs=10, bold=True, color=fc)
        tb(slide, grp,  col_xs[1], ry+In(0.1), col_ws[1], rh-In(0.15), fn='Calibri', fs=9, color=AM if is_ms else G)
        tb(slide, mas,  col_xs[2], ry+In(0.1), col_ws[2], rh-In(0.15), fn='Courier New', fs=9, color=AM if is_ms else T2)
        score_str = f'{e}/{total} exact\n{a}/{total} adj'
        bg_score = bc; bc_score = fc
        rect(slide, col_xs[3]+In(0.05), ry+In(0.06), col_ws[3]-In(0.1), rh-In(0.14), fill=bg_score, line=bc_score, lw=Pt(0.75))
        tb(slide, score_str, col_xs[3]+In(0.12), ry+In(0.08), col_ws[3]-In(0.2), rh-In(0.18),
           fn='Courier New', fs=9, bold=True, color=fc)
        tb(slide, note, col_xs[4], ry+In(0.06), col_ws[4]-In(0.1), rh-In(0.12),
           fn='Calibri', fs=8.5, italic=True, color=T3)
        ry += rh + In(0.04)


def s9_param_bias(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'Parameter Accuracy')
    title(slide, 'Swing Count & First-Angle Bias vs OptiTrack', top=In(0.50))
    hline(slide, In(1.08))
    tb(slide, 'Positive = overestimates  ·  Negative = underestimates  ·  Zero = perfect agreement',
       In(0.55), In(1.12), In(8.9), In(0.22), fn='Calibri', fs=9.5, italic=True, color=T3)

    # N Bias table
    tb(slide, 'N BIAS — SWING COUNT ERROR (cycles)', In(0.55), In(1.38), In(4.35), In(0.22),
       fn='Courier New', fs=8, bold=True, color=T3)
    n_data = [
        ('MediaPipe Full',  '+4.03', AM),
        ('MediaPipe Lite',  '+3.68', AM),
        ('pose2sim_lw',     '+2.69', AM),
        ('RTMPose-S',       '+1.78', AM),
        ('MediaPipe Heavy', '+1.48', AM),
        ('HRNet-W48',       '+1.33', AM),
        ('RTMPose-M',       '+1.06', G),  # best (highlighted)
        ('DETRPose-S',      '-2.71', BL),
        ('YOLO-N',          '-3.33', RD),
    ]
    bar_l = In(1.95); bar_zero = In(0.5); bar_scale = In(0.55)
    ry = In(1.62)
    for lbl, val, col in n_data:
        v = float(val)
        is_best = col == G
        bh = In(0.27)
        rect(slide, In(0.55), ry, In(4.35), bh, fill=GB if is_best else WH, line=BR, lw=Pt(0.3))
        tb(slide, lbl, In(0.6), ry+In(0.04), In(1.3), bh-In(0.06),
           fn='Courier New', fs=8.5, bold=is_best, color=G if is_best else T2)
        bw = abs(v) * bar_scale
        bx = bar_l + bar_zero + (0 if v < 0 else 0)
        if v >= 0:
            bx = bar_l + bar_zero
        else:
            bx = bar_l + bar_zero - bw
        rect(slide, bx, ry+In(0.04), max(bw, Pt(2)), bh-In(0.1), fill=col)
        tb(slide, val, bx + (bw+In(0.04) if v>=0 else -In(0.44)), ry+In(0.05), In(0.42), bh-In(0.1),
           fn='Courier New', fs=8.5, bold=is_best, color=col)
        ry += bh + In(0.02)

    # A0 Bias table
    tb(slide, 'A₀ BIAS — FIRST-ANGLE ERROR (degrees)', In(4.65), In(1.38), In(4.7), In(0.22),
       fn='Courier New', fs=8, bold=True, color=T3)
    a0_data = [
        ('RTMPose-S',       '-14.42', BL),
        ('DETRPose-S',      '-11.23', BL),
        ('RTMPose-M',        '-9.69', BL),
        ('MediaPipe Heavy',  '-9.59', BL),
        ('MediaPipe Lite',   '-9.19', BL),
        ('MediaPipe Full',   '-8.48', BL),
        ('pose2sim_lw',      '-5.45', BL),
        ('YOLO-N',           '-4.28', BL),
        ('HRNet-W48',        '+6.35', RD),  # only overestimator — highlighted
    ]
    rx = In(4.65); r_bar_zero = In(1.7); a0_scale = In(0.1)
    ry2 = In(1.62)
    for lbl, val, col in a0_data:
        v = float(val)
        is_bad = col == RD
        bh = In(0.27)
        rect(slide, rx, ry2, In(4.7), bh, fill=RB if is_bad else WH, line=BR, lw=Pt(0.3))
        tb(slide, lbl, rx+In(0.05), ry2+In(0.04), In(1.3), bh-In(0.06),
           fn='Courier New', fs=8.5, bold=is_bad, color=RD if is_bad else T2)
        bw = abs(v) * a0_scale
        bx = rx + r_bar_zero + (0 if v < 0 else 0)
        if v >= 0:
            bx = rx + r_bar_zero
        else:
            bx = rx + r_bar_zero - bw
        rect(slide, bx, ry2+In(0.04), max(bw, Pt(2)), bh-In(0.1), fill=col)
        tb(slide, val+'°', bx+(bw+In(0.04) if v>=0 else -In(0.52)), ry2+In(0.05), In(0.52), bh-In(0.1),
           fn='Courier New', fs=8.5, bold=is_bad, color=col)
        ry2 += bh + In(0.02)

    # Footer note
    rect(slide, In(0.55), In(5.0), In(8.9), In(0.5), fill=PN, line=BR)
    tb(slide, 'RTMPose-M has the smallest N bias (+1.06 cycles). All MediaPipe variants overcount by 1.5–4 swings. '
              'HRNet-W48 is the ONLY model that overestimates A₀ (+6.35°) — all others underestimate by 4°–14°, '
              'compressing PT scores toward healthier-looking grades. HRNet\'s A₀ overestimate pushes it to overcall spasticity in healthy participants.',
       In(0.7), In(5.04), In(8.6), In(0.42), fn='Calibri', fs=9, color=T2)


def s10_findings(prs):
    slide = blank_slide(prs)
    eyebrow(slide, 'Summary')
    title(slide, 'Key Findings', top=In(0.50))
    hline(slide, In(1.08))

    findings = [
        (PV, PB, '1', 'VIEWER',
         'Pendulastic Viewer ranked #2 among all 11 systems',
         '85.3% adjacent MAS agreement (34 sessions). Outperforms MediaPipe Full (84.4%) and Lite (83.3%). '
         'Stronger PT correlation (r=0.394 vs 0.156). Delivers clinical-grade MAS estimation from an ordinary webcam.'),
        (AM, AB, '2', 'GAP',
         'Viewer underdetects MS spasticity — 53.3% vs 87.5% for MediaPipe Full',
         'Only 8/15 MS sessions scored MAS ≥1. Partly explained by the P5 paradox. '
         'Opportunity: MS-specific threshold calibration may close this gap.'),
        (AM, AB, '3', 'DISCOVERY',
         'P5 paradox: 3D biomechanics and clinical assessment diverge',
         'P5 is an MS patient, yet OptiTrack scores every trial MAS 0. Viewer scores MAS 1–3 on 6/11 sessions — '
         'possibly detecting real subclinical spasticity that pendulum kinematics miss. Warrants clinical review.'),
        (G,  GB, '4', 'BEST HPE',
         'Pose2Sim-Lightweight leads HPE at 92.3% adj — needs more coverage',
         'Strongest HPE model: r=0.773 PT correlation, 92.3% adjacent agreement. '
         'But n=13 trials limits confidence. Priority: run Pose2Sim on the full cohort.'),
        (RD, RB, '5', 'DISCARD',
         'YOLO-N (6.2%) and DETRPose-S (37.9%) not viable for PT scoring',
         'YOLO-N: extreme swing undercounting, MAE=0.578. DETRPose-S achieves MS recall only by '
         'always scoring non-zero — clinically uninformative. Neither should be used for PT assessment.'),
    ]
    fh = In(0.82)
    fy = In(1.18)
    for fc, bc, num, tag, head, desc in findings:
        rect(slide, In(0.55), fy, In(8.9), fh, fill=WH, line=BR, lw=Pt(0.5))
        rect(slide, In(0.55), fy, Pt(4), fh, fill=fc)
        # Number
        tb(slide, num, In(0.7), fy+In(0.15), In(0.3), In(0.4), fn='Georgia', fs=22, color=BR)
        # Tag badge
        rect(slide, In(1.1), fy+In(0.15), In(0.7), In(0.22), fill=bc)
        tb(slide, tag, In(1.1), fy+In(0.16), In(0.7), In(0.2),
           fn='Courier New', fs=7.5, bold=True, color=fc, align=PP_ALIGN.CENTER)
        # Heading
        tb(slide, head, In(1.88), fy+In(0.1), In(7.5), In(0.28), fn='Georgia', fs=13, color=T1)
        # Description
        tb(slide, desc, In(1.88), fy+In(0.4), In(7.5), In(0.38), fn='Calibri', fs=9, color=T2)
        fy += fh + In(0.04)


def make_waveform_figure():
    """Return a BytesIO PNG of an annotated healthy vs MS pendulum waveform."""

    # ── Palette matching slide ──────────────────────────────────────────────
    Cg  = '#0A7B52'  # green  – healthy / good
    Ca  = '#B06010'  # amber  – MS / A0 annotation
    Cb  = '#1565C0'  # blue   – A1 / R2n
    Cp  = '#6030A0'  # purple – A2 / phi_max
    Cr  = '#C03030'  # red    – omega_max / N
    Cbg = '#F0F5FA'  # slide background
    Cl  = '#B8CDDE'  # light border / neutral line
    Cm  = '#7A9CB8'  # muted labels

    # ── Model parameters ────────────────────────────────────────────────────
    # Calibrated to match study medians: healthy N=7, phi_max=0.79, R2n=1.01
    A0_h, f_h, lam_h = 65.0, 0.90, 0.25
    # MS: N≈2, phi_max≈0.19 (representative of P4-right severity)
    A0_m, f_m, lam_m = 57.0, 0.85, 1.50
    T_h, T_m = 1.0 / f_h, 1.0 / f_m

    t = np.linspace(0.0, 9.0, 6000)
    th_h = A0_h * np.exp(-lam_h * t) * np.cos(2 * np.pi * f_h * t)
    th_m = A0_m * np.exp(-lam_m * t) * np.cos(2 * np.pi * f_m * t)

    # ── Derived annotated values (healthy) ──────────────────────────────────
    h_trough  = A0_h * np.exp(-lam_h * T_h / 2) * np.cos(np.pi)   # first flexion trough (neg)
    h_A2      = A0_h * np.exp(-lam_h * T_h)                         # first return peak (pos)
    h_A1pp    = A0_h + abs(h_trough)                                 # peak-to-peak of first osc
    h_phi     = h_A2 / A0_h
    h_R2n     = h_A1pp / (1.6 * A0_h)
    # omega_max: numerical max |velocity| in first half-cycle
    vel_h     = np.gradient(th_h, t[1] - t[0])
    early     = t < T_h / 2
    idx_om    = np.argmax(np.abs(vel_h[early]))
    t_om, th_om, slope_om = t[early][idx_om], th_h[early][idx_om], vel_h[early][idx_om]
    # N: count cycles until envelope < 1.5°
    env_h       = A0_h * np.exp(-lam_h * t)
    si_h        = np.argmax(env_h < 1.5)
    t_settle_h  = t[si_h] if si_h > 0 else 9.0
    h_N         = max(1, round(f_h * t_settle_h))

    # ── Derived annotated values (MS) ───────────────────────────────────────
    m_A2    = A0_m * np.exp(-lam_m * T_m)
    m_phi   = m_A2 / A0_m
    env_m   = A0_m * np.exp(-lam_m * t)
    si_m    = np.argmax(env_m < 1.5)
    t_settle_m = t[si_m] if si_m > 0 else 9.0
    m_N     = max(1, round(f_m * t_settle_m))

    # ── Figure ──────────────────────────────────────────────────────────────
    fig = plt.figure(figsize=(11.5, 3.85), facecolor=Cbg)
    ax1 = fig.add_axes([0.085, 0.20, 0.545, 0.70])   # healthy (wide)
    ax2 = fig.add_axes([0.690, 0.20, 0.285, 0.70])   # MS (narrow)

    def _style(ax):
        ax.set_facecolor(Cbg)
        for side, vis in [('top', False), ('right', False),
                          ('left', True),  ('bottom', True)]:
            ax.spines[side].set_visible(vis)
            if vis:
                ax.spines[side].set_color(Cl)
                ax.spines[side].set_linewidth(0.8)
        ax.tick_params(colors=Cm, labelsize=7.5, length=3, width=0.7)
        ax.set_xlabel('Time (s)', color=Cm, fontsize=8)
        ax.set_ylabel('Knee angle from neutral (°)', color=Cm, fontsize=8)
        ax.axhline(0, color=Cl, lw=1.0, ls='--', zorder=0)
        ax.set_yticks([-60, -30, 0, 30, 60])
        ax.set_ylim(-92, 108)

    _style(ax1); _style(ax2)

    aw = dict(arrowstyle='<->', mutation_scale=10, lw=1.4)   # double-head arrow kwargs

    # ────────────────────────── Healthy panel ─────────────────────────────────
    ax1.set_xlim(-0.65, 9.0)
    ax1.set_title('Healthy Control — freely oscillating pendulum',
                  color=Cg, fontsize=9, fontweight='bold', loc='left', pad=7)
    ax1.plot(t, th_h, color=Cg, lw=2.0, zorder=5)

    # Shade areas above/below neutral for area_ratio
    ax1.fill_between(t, 0, th_h, where=(th_h > 0), color=Cg, alpha=0.08, zorder=1)
    ax1.fill_between(t, 0, th_h, where=(th_h < 0), color=Ca, alpha=0.08, zorder=1)
    ax1.text(6.5,  18, 'area (+)', color=Cg, fontsize=6.5, alpha=0.8)
    ax1.text(6.5, -22, 'area (−)', color=Ca, fontsize=6.5, alpha=0.8)
    ax1.text(6.5, -35, f'area_ratio ≈ 0.09\n(balanced → healthy)', color=Cm, fontsize=6.5)

    # ── A0 ──
    ax1.annotate('', xy=(0, A0_h), xytext=(0, 0),
                 arrowprops=dict(**aw, color=Ca))
    ax1.text(-0.47, A0_h / 2,
             f'A₀\n{A0_h:.0f}°\n(drop)', color=Ca,
             fontsize=7.5, ha='center', va='center', fontweight='bold')

    # ── First trough (for A1 context) ──
    ax1.plot(T_h / 2, h_trough, 'o', color=Cb, ms=4.5, zorder=7)
    ax1.annotate('', xy=(T_h / 2, h_trough), xytext=(T_h / 2, 0),
                 arrowprops=dict(**aw, color=Cb))
    ax1.text(T_h / 2 + 0.07, h_trough * 0.55,
             f'{abs(h_trough):.0f}°', color=Cb, fontsize=7.5, va='center')

    # ── A1 peak-to-peak bracket (left of trough) ──
    x_br = T_h / 2 - 0.20
    ax1.annotate('', xy=(x_br, A0_h), xytext=(x_br, h_trough),
                 arrowprops=dict(arrowstyle='<->', mutation_scale=10, lw=1.1, color=Cb))
    ax1.text(x_br - 0.08, (A0_h + h_trough) / 2,
             f'A₁={h_A1pp:.0f}°\n(pk-pk)',
             color=Cb, fontsize=7, ha='right', va='center',
             bbox=dict(boxstyle='round,pad=0.2', fc='white', ec=Cb, lw=0.7, alpha=0.93))
    # R2n formula below waveform
    ax1.text(0.25, -78,
             f'R2n = A₁ / (1.6 × A₀)  =  {h_A1pp:.0f} / {1.6*A0_h:.0f}  =  {h_R2n:.2f}',
             color=Cb, fontsize=7.5,
             bbox=dict(boxstyle='round,pad=0.25', fc='white', ec=Cb, lw=0.7, alpha=0.93))

    # ── A2 and phi_max ──
    ax1.plot(T_h, h_A2, 'o', color=Cp, ms=5.5, zorder=7)
    ax1.annotate('', xy=(T_h, h_A2), xytext=(T_h, 0),
                 arrowprops=dict(**aw, color=Cp))
    ax1.text(T_h + 0.09, h_A2 / 2,
             f'A₂\n{h_A2:.0f}°', color=Cp, fontsize=7.5, va='center', fontweight='bold')
    ax1.text(T_h + 0.75, 87,
             f'phi_max = A₂ / A₀\n= {h_A2:.0f} / {A0_h:.0f}  =  {h_phi:.2f}',
             color=Cp, fontsize=8.5, fontweight='bold', va='top',
             bbox=dict(boxstyle='round,pad=0.3', fc='white', ec=Cp, lw=1.2, alpha=0.93))

    # ── omega_max tangent ──
    dt_t = 0.17
    t_tang = np.array([t_om - dt_t, t_om + dt_t])
    th_tang = th_om + slope_om * (t_tang - t_om)
    ax1.plot(t_tang, th_tang, color=Cr, lw=2.2, zorder=8)
    ax1.text(t_om + 0.20, th_om - 9,
             f'ω_max = {abs(slope_om):.0f}°/s\n'
             f'ω_max_n = ω_max / A₀ = {abs(slope_om)/A0_h:.2f} /s',
             color=Cr, fontsize=7.5, va='top',
             bbox=dict(boxstyle='round,pad=0.25', fc='white', ec=Cr, lw=0.7, alpha=0.93))

    # ── Period / frequency bracket ──
    y_per = -85
    ax1.annotate('', xy=(T_h, y_per), xytext=(0.0, y_per),
                 arrowprops=dict(**aw, color=Cm))
    ax1.text(T_h / 2, y_per - 7,
             f'T = {T_h:.2f} s  →  f = {f_h:.2f} Hz',
             color=Cm, fontsize=7.5, ha='center', va='top')

    # ── N swings bracket ──
    y_N = 96
    ax1.annotate('', xy=(t_settle_h, y_N), xytext=(0.02, y_N),
                 arrowprops=dict(arrowstyle='->', color=Cg, lw=1.3, mutation_scale=10))
    ax1.text(t_settle_h / 2, y_N + 4,
             f'N = {h_N} complete swings  (oscillations > 1°)',
             color=Cg, fontsize=8.5, ha='center', va='bottom', fontweight='bold')

    # Neutral label
    ax1.text(8.8, 2.5, 'neutral', color=Cm, fontsize=7, ha='right')

    # ────────────────────────── MS panel ─────────────────────────────────────
    ax2.set_xlim(-0.65, 9.0)
    ax2.set_title('MS Patient — spastic resistance',
                  color=Ca, fontsize=9, fontweight='bold', loc='left', pad=7)
    ax2.plot(t, th_m, color=Ca, lw=2.0, zorder=5)

    # A0
    ax2.annotate('', xy=(0, A0_m), xytext=(0, 0),
                 arrowprops=dict(**aw, color=Ca))
    ax2.text(-0.47, A0_m / 2, f'A₀\n{A0_m:.0f}°',
             color=Ca, fontsize=7.5, ha='center', va='center', fontweight='bold')

    # A2 (tiny)
    ax2.plot(T_m, m_A2, 'o', color=Cp, ms=5.5, zorder=7)
    ax2.annotate('', xy=(T_m, max(m_A2, 1.5)), xytext=(T_m, 0),
                 arrowprops=dict(**aw, color=Cp))
    ax2.text(T_m + 0.1, max(m_A2, 2) / 2 + 3,
             f'A₂ = {m_A2:.0f}°', color=Cp, fontsize=7.5, va='center', fontweight='bold')

    ax2.text(1.6, 88,
             f'phi_max = {m_phi:.2f}\n(healthy: 0.79)',
             color=Cp, fontsize=9, fontweight='bold', va='top',
             bbox=dict(boxstyle='round,pad=0.3', fc='white', ec=Cp, lw=1.5, alpha=0.93))
    ax2.text(1.6, 55,
             f'N = {m_N} swings\n(healthy: {h_N})',
             color=Cr, fontsize=9, fontweight='bold', va='top',
             bbox=dict(boxstyle='round,pad=0.3', fc='white', ec=Cr, lw=1.5, alpha=0.93))

    # Spastic catch annotation
    ax2.annotate('spastic\ncatch',
                 xy=(t_settle_m, env_m[si_m if si_m > 0 else -1]),
                 xytext=(t_settle_m + 1.2, 28),
                 color=Ca, fontsize=7.5, style='italic',
                 arrowprops=dict(arrowstyle='->', color=Ca, lw=1.0, mutation_scale=8))

    ax2.text(8.8, 2.5, 'neutral', color=Cm, fontsize=7, ha='right')

    # ── Vertical separator between panels ──
    fig.add_artist(plt.Line2D([0.645, 0.645], [0.05, 0.98],
                               transform=fig.transFigure, color=Cl, lw=1.0))

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor=Cbg, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def s_waveform(prs):
    """Slide: annotated pendulum test waveform showing where each parameter comes from."""
    slide = blank_slide(prs)
    eyebrow(slide, 'Pendulum Test — Waveform Anatomy')
    tb(slide, 'Reading the Angle Waveform: Where Each Parameter Comes From',
       In(0.55), In(0.44), In(9.0), In(0.42),
       fn='Georgia', fs=20, color=T1)
    tb(slide,
       'Synthetic waveforms calibrated to study medians  ·  '
       'Knee angle (degrees) vs time after leg is released from full extension  ·  '
       'Left: Healthy control (N=7, phi_max=0.79)  ·  Right: MS patient (N=2, phi_max=0.19)',
       In(0.55), In(0.90), In(9.0), In(0.22),
       fn='Calibri', fs=8.5, italic=True, color=T3)
    hline(slide, In(1.14))

    buf = make_waveform_figure()
    slide.shapes.add_picture(buf, In(0.40), In(1.17), In(9.25), In(4.30))


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    print("Building slides...")
    s1_title(prs)
    print("  1/9 Title")
    s2_study_design(prs)
    print("  2/9 Study Design")
    s3_optitrack(prs)
    print("  3/9 How OptiTrack Works")
    s4_hpe_tracking(prs)
    print("  4/9 How HPE Models Track the Knee")
    s5_pendulastic_algorithm(prs)
    print("  5/11 Pendulastic Algorithm")
    s_waveform(prs)
    print("  6/11 Waveform Anatomy")
    s_param_profile(prs)
    print("  7/11 Parameter Profiles: MS vs Control")
    s6_viewer_results(prs)
    print("  7/10 Viewer Results")
    s7_rankings(prs)
    print("  7/9 Full Rankings")
    s8_per_participant(prs)
    print("  8/9 Per-Participant")
    s9_param_bias(prs)
    print("  9/9 Parameter Bias")
    s10_findings(prs)
    print(" 10/9 Key Findings")

    out = os.path.join(os.path.dirname(__file__),
                       "Pendulastic_Benchmark.pptx")
    prs.save(out)
    print(f"\nSaved -> {out}")
    return out

if __name__ == "__main__":
    main()
