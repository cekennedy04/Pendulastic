"""
build_results_pptx.py
======================
Builds a publication-style PowerPoint deck for the Results section (2026-08-20),
mirroring the Google Doc text delivered the same day. Reuses the already-generated
figures in Model_Analysis_Outputs/paper_figures/ -- no new analysis, presentation only.

Usage:
    .venv\\Scripts\\python.exe build_results_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

FIG_DIR = "Model_Analysis_Outputs/paper_figures"
OUT_PATH = "Model_Analysis_Outputs/Pendulastic_Results_Section.pptx"

INK = RGBColor(0x1a, 0x1a, 0x1a)
MUTED = RGBColor(0x6b, 0x6b, 0x6b)
BLUE = RGBColor(0x2a, 0x78, 0xd6)
ORANGE = RGBColor(0xeb, 0x68, 0x34)
WHITE = RGBColor(0xff, 0xff, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(12.83), Inches(1.3))
    line.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)
    line.line.width = Pt(1)


def add_bullets(slide, items, left=0.5, top=1.5, width=12.3, height=5.3, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(10)
    return box


def add_figure(slide, filename, left, top, width):
    path = os.path.join(FIG_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    else:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1))
        box.text_frame.text = f"[missing: {filename}]"


def add_table(slide, headers, rows, left, top, width, height, font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].font.color.rgb = INK
    return table


# ---------------------------------------------------------------- Slide 1: Title
s = add_slide()
box = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Smartphone-IMU Validation of the Wartenberg Pendulum Test for Spasticity Assessment"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = INK
box2 = s.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.6))
p2 = box2.text_frame.paragraphs[0]
p2.text = "Results"
p2.font.size = Pt(20)
p2.font.color.rgb = ORANGE
box3 = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
p3 = box3.text_frame.paragraphs[0]
p3.text = "Pendulastic  |  Prepared for review with Dr. Monica Perez  |  2026-08-20"
p3.font.size = Pt(13)
p3.font.color.rgb = MUTED

# ---------------------------------------------------------------- Slide 2: Sample
s = add_slide()
add_title(s, "3.1  Participant Characteristics and Data Availability")
add_bullets(s, [
    "14 enrolled participants: 8 healthy control, 6 MS  (+1 additional MS participant, P17, grading in progress -- excluded from clinical comparisons)",
    "Only 5 participants (4 MS: P5, P13, P14, P15; 1 control: P16) have synchronized IMU + OptiTrack data -- this is the binding constraint on every result in this deck",
    "61 trials: IMU + OptiTrack matched sample",
    "37 trials: video + OptiTrack matched sample",
    "49 trials: three-way matched (IMU + video + OptiTrack) -- used for Sections 3.7-3.8",
    "40 trials: complete production PT-score parameter sets -- used for Section 3.9",
    "MAS grades available for all 61 IMU trials; current cohort tops out at MAS = 1+ (no MAS ≥ 2 trials yet)",
])

# ---------------------------------------------------------------- Slide 3: IMU agreement
s = add_slide()
add_title(s, "3.2  IMU vs. OptiTrack Agreement", "n = 61 trials, 5 participants")
add_table(s, ["Parameter", "ICC(2,1)", "Bias", "95% LoA"], [
    ["R2n (relaxation index)", "0.226", "−0.108", "[−1.06, 0.85]"],
    ["N (oscillation count)", "0.458", "−1.46", "[−7.48, 4.56]"],
    ["φmax ratio", "0.044", "−0.055", "[−0.51, 0.40]"],
    ["ωmax,n", "0.014", "−2.13", "[−12.05, 7.78]"],
    ["f (frequency)", "0.140", "−0.40", "[−1.98, 1.18]"],
    ["Area ratio", "0.135", "−0.052", "[−0.82, 0.72]"],
    ["ωmin,n", "0.214", "−1.01", "[−6.70, 4.69]"],
], left=0.5, top=1.5, width=6.2, height=3.6, font_size=11)
add_figure(s, "fig1_bland_altman.png", 7.0, 1.5, 5.8)
box = s.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Trajectory-level RMSE: 14.84° mean / 10.98° median (n=53); only 2/53 trials met "
          "the 5° clinical-goal threshold. 67–76% of total error is a fixed per-trial "
          "calibration bias, not random noise -- the primary target for future correction.")
p.font.size = Pt(14)
p.font.color.rgb = INK

# ---------------------------------------------------------------- Slide 4: MediaPipe agreement
s = add_slide()
add_title(s, "3.3  MediaPipe vs. OptiTrack Agreement", "n = 49 trials, 5 participants")
add_table(s, ["Parameter", "ICC(2,1)", "Bias", "95% LoA"], [
    ["R2n (relaxation index)", "−0.041", "−0.441", "[−2.35, 1.47]"],
    ["N (oscillation count)", "0.032", "−2.30", "[−18.35, 13.75]"],
    ["φmax ratio", "−0.008", "0.110", "[−0.91, 1.13]"],
    ["ωmax,n", "−0.036", "6.75", "[−50.28, 63.77]"],
    ["f (frequency)", "−0.115", "−0.83", "[−3.04, 1.37]"],
    ["Area ratio", "0.003", "−0.420", "[−1.23, 0.39]"],
    ["ωmin,n", "−0.018", "3.69", "[−29.16, 36.53]"],
], left=0.5, top=1.5, width=6.2, height=3.6, font_size=11)
add_figure(s, "fig31_mp_trajectory_example.png", 7.0, 1.5, 5.8)
box = s.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Full-curve RMSE 36.0° mean / 33.3° median (n=37); 0/37 met the 5° threshold. "
          "ICC ≤ 0 for all 7 parameters -- no measurable agreement, not merely \"poor.\" "
          "Right: representative single-trial trajectory (MediaPipe vs. OptiTrack) showing the tracking failure.")
p.font.size = Pt(14)
p.font.color.rgb = INK

# ---------------------------------------------------------------- Slide 5: Clinical validity
s = add_slide()
add_title(s, "3.4  Clinical Validity: Correlation With MAS")
add_bullets(s, [
    "IMU-derived relaxation index (R2n) vs. MAS grade:",
    "     Spearman's ρ = −0.313,  p = 0.014  (n = 61 trials, all 5 participants)",
    "Direction is as expected: lower R2n (more restricted swing) ↔ higher MAS grade",
    "Real and statistically significant -- but modest relative to the closest comparable",
    "     published result: Yeh et al. (2025), ρ = −0.75 to −0.78 (stroke, n = 20,",
    "     wider severity range than currently available here)",
])

# ---------------------------------------------------------------- Slide 6: Group comparison
s = add_slide()
add_title(s, "3.5  Group Comparison: MS vs. Control")
add_figure(s, "fig4_metrics_by_group.png", 0.6, 1.5, 7.2)
add_bullets(s, [
    "Linear mixed-effects model (R2n ~ group, participant random intercept)",
    "MS estimated 0.173 lower than control (expected direction)",
    "β = −0.173, SE = 0.207, z = −0.835, p = 0.404",
    "Not significant -- but NOT a null finding:",
    "control arm = 1 participant (P16)",
    "This compares 4 MS people to one individual,",
    "not to a control distribution.",
], left=8.0, top=1.6, width=4.8, size=15)

# ---------------------------------------------------------------- Slide 7: Pre/post
s = add_slide()
add_title(s, "3.6  Longitudinal Change (Illustrative Only)")
add_figure(s, "fig5_pre_post.png", 0.6, 1.5, 7.2)
add_bullets(s, [
    "Only P15 has both pre- and post-treatment",
    "recordings in the OptiTrack-matched dataset",
    "R2n: 0.91 → 0.94 (pre → post)",
    "Direction consistent with reduced spasticity",
    "",
    "This is one paired observation.",
    "No statistical inference is drawn.",
    "Not evidence of treatment effect.",
], left=8.0, top=1.6, width=4.8, size=15)

# ---------------------------------------------------------------- Slide 8: sign inversion
s = add_slide()
add_title(s, "3.7  Do Parameters Discriminate Consistently Across Modalities?")
add_figure(s, "fig6_metric_effect_heatmap.png", 0.6, 1.4, 5.6)
add_bullets(s, [
    "Cohen's d (MAS>0 vs. MAS=0), 7 parameters × 3 modalities,",
    "49 three-way-matched trials",
    "",
    "OptiTrack (ground truth): positive effect sizes for",
    "φmax ratio, ωmax,n, f",
    "",
    "IMU and MediaPipe: OPPOSITE SIGN on nearly all",
    "of the same parameters, same trials",
    "",
    "More serious than noisy agreement: a single-modality",
    "parameter can point the wrong clinical direction,",
    "not just be imprecise around the right one.",
], left=6.6, top=1.4, width=6.2, size=15)

# ---------------------------------------------------------------- Slide 9: AUC
s = add_slide()
add_title(s, "3.8  Single-Metric vs. Combined-Metric Classification")
add_figure(s, "fig7_single_vs_combined_auc.png", 0.6, 1.4, 7.6)
add_bullets(s, [
    "Leave-one-participant-out cross-validated",
    "logistic regression (the honest generalization test)",
    "",
    "n = 49 trials, 5 participants",
    "(35 MAS>0 / 14 MAS=0)",
    "",
    "Every single-parameter AUC < 0.5",
    "(range 0.0–0.17)",
    "",
    "Combined 7-parameter AUC = 0.21",
    "-- worse than chance",
    "",
    "Reflects sample size (5 people),",
    "not refutation of the hypothesis.",
], left=8.4, top=1.4, width=4.4, size=14)

# ---------------------------------------------------------------- Slide 10: composite score
s = add_slide()
add_title(s, "3.9  Does the Production Composite Score Separate Groups?")
add_figure(s, "fig8_score_naive_vs_logocv.png", 0.6, 1.4, 7.6)
add_bullets(s, [
    "compute_pt_score() run on 40 IMU trials,",
    "5 participants (complete parameter sets)",
    "",
    "Naive Mann-Whitney U (Control vs. MS):",
    "p = 0.5865 -- NOT significant",
    "",
    "Contradicts prior internal documentation:",
    "\"p = 0.0001\" (earlier PT7 comparison)",
    "-- flagged for re-verification",
    "",
    "LOGO-CV version could not even be run:",
    "only 1 control participant in matched set",
], left=8.4, top=1.4, width=4.4, size=14)

# ---------------------------------------------------------------- Slide 11: collinearity
s = add_slide()
add_title(s, "3.10  Internal Structure: Parameter Collinearity")
add_figure(s, "fig10_param_correlation.png", 0.6, 1.4, 6.8)
add_bullets(s, [
    "All 7 parameters weighted EQUALLY (1/7 each)",
    "-- not fit or validated against outcomes",
    "",
    "R2n and ωmax,n correlate at r = 0.93",
    "-- near-redundant, not independent",
    "",
    "Shared signal effectively double-counted;",
    "distinct parameters (e.g. area ratio) under-weighted",
    "",
    "HEALTHY_REF anchored on only 4 controls",
    "-- flagged \"PROVISIONAL\" in the code itself",
    "-- leave-one-control-out check could not run",
    "  (only 1 control in the matched set)",
], left=7.6, top=1.4, width=5.2, size=14)

# ---------------------------------------------------------------- Slide 12: MAS target
s = add_slide()
add_title(s, "3.11  Is MAS Correlation Measuring the Right Clinical Target?")
add_bullets(s, [
    "The pendulum test targets knee-EXTENSOR spasticity specifically -- not spasticity generally",
    "Section 3.4's correlation used mas_grade, a collapsed field (flexor + extensor folded together)",
    "",
    "For P15, the one participant with granular data available:",
    "     mas_grade:      0 (left)  /  1 (right)",
    "     mas_flexion:    1+ (left) /  1 (right)",
    "     mas_extension:  0 (left)  /  0 (right)  ← zero on both legs",
    "",
    "Open question this dataset cannot yet resolve: does R2n track the extensor-specific",
    "mechanism the test is designed for, or flexor signal that happens to covary with the",
    "collapsed grade in this small sample?",
    "",
    "A newly enrolled participant (P17) has flexion/extension scored before an overall grade",
    "was assigned -- suggesting the sub-components may be the more appropriate target going forward.",
], size=15)

# ---------------------------------------------------------------- Slide 13: Summary table
s = add_slide()
add_title(s, "Summary of Findings")
add_table(s, ["Question", "Result"], [
    ["IMU agree with OptiTrack?", "Modest (ICC 0.01–0.46); ~70% of error is correctable bias"],
    ["MediaPipe agree with OptiTrack?", "No measurable agreement (ICC ≤ 0)"],
    ["IMU index correlate with MAS?", "Yes, significant (ρ=−0.313, p=0.014); weaker than published comparator"],
    ["MS vs. control differ on parameters?", "Expected direction, n.s.; uninterpretable with n=1 control"],
    ["Combining parameters improve classification?", "Not yet demonstrable -- all AUCs below chance, n=5 participants"],
    ["Composite score separate Control/MS?", "No, even naively (p=0.5865); contradicts prior p=0.0001 claim"],
    ["Composite score internally sound?", "No -- r=0.93 collinear pair; n=4 healthy reference untested"],
    ["MAS correlation measuring right target?", "Unresolved -- test is extensor-specific; only granular case shows 0 extensor"],
], left=0.5, top=1.5, width=12.3, height=4.8, font_size=13)
box = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("All limitations trace to one addressable cause: only 5 of 15 enrolled participants have "
          "synchronized IMU + OptiTrack data. IMU-recording the 7 existing video-only controls -- zero "
          "new recruitment required -- is the single highest-value next step.")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = MUTED

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH}")
